#!/usr/bin/env python3
"""preflight_check must catch the mechanical defects AND refuse to imply it covered the rest.

Script-style (main() + explicit exit) to match the other suites; pytest collects nothing here
and ci.yml invokes it directly.

The second half of this suite matters as much as the first: a checker that quietly reported
"clean" for a judgment item it never examined would be worse than no checker, because the
twelve ticks would then LOOK mechanised while five of them were still self-attested.
"""
from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SKILL = HERE.parent
TOOL = SKILL / "scripts" / "preflight_check.py"
sys.path.insert(0, str(SKILL / "scripts"))


def build(path: Path, *, meta=False, notes=True, asof=True):
    import deckkit as dk
    from pptx.dml.color import RGBColor
    C = lambda h: RGBColor.from_string(h)
    prs = dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dk.text(s, 0.6, 1.0, 8, 0.6,
            [[("Results (editable native chart)" if meta else "Agents act, not answer",
               24, C("12395E"), True, False, "Helvetica Neue")]])
    if asof:
        dk.text(s, 0.6, 4.8, 8, 0.3,
                [[("Figures as of 28 July 2026", 11, C("646F7B"), False, False, "Helvetica Neue")]])
    if notes:
        dk.speaker_notes(s, "spoken thread")
    prs.save(str(path))
    return path


def run(deck, *extra):
    p = subprocess.run([sys.executable, str(TOOL), str(deck), *extra],
                       capture_output=True, text=True)
    return p.returncode, p.stdout + p.stderr


CASES = []


def case(name):
    def deco(fn):
        CASES.append((name, fn))
        return fn
    return deco


@case("leaked meta-annotation on a slide is caught")
def _(td):
    code, out = run(build(td / "meta.pptx", meta=True), "--static")
    return code == 1 and "editable native chart" in out


@case("unfilled <slot> template text is caught")
def _(td):
    import deckkit as dk
    from pptx.dml.color import RGBColor
    prs = dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dk.text(s, 0.6, 1.0, 8, 0.6, [[("Cost was <amount> per run", 20,
            RGBColor.from_string("12395E"), False, False, "Helvetica Neue")]])
    dk.speaker_notes(s, "x")
    prs.save(str(td / "slot.pptx"))
    code, out = run(td / "slot.pptx", "--static")
    return code == 1 and "<amount>" in out


@case("missing speaker notes on a presented deck is caught")
def _(td):
    code, out = run(build(td / "nonotes.pptx", notes=False), "--static")
    return code == 1 and "NO speaker notes" in out


@case("missing as-of date is reported but does NOT fail the run")
def _(td):
    # Advisory on purpose: a deck with no time-bound claims legitimately carries no date, and
    # nothing in the file distinguishes the two. Failing here would train people to ignore the
    # tool, which costs more than the miss.
    code, out = run(build(td / "nodate.pptx", asof=False), "--static")
    return code == 0 and "as of" in out.lower() and "advisory" in out.lower()


@case("a font missing on THIS machine is advisory, not a failure")
def _(td):
    # The regression that broke CI: fixtures used Helvetica Neue, absent on Ubuntu, so item 10
    # failed and every "expects a clean run" case failed with it. The deck was fine; the check
    # was wrong. Uses a font guaranteed absent everywhere so the assertion is platform-stable.
    import deckkit as dk
    from pptx.dml.color import RGBColor
    prs = dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dk.text(s, 0.6, 1.0, 8, 0.6, [[("Title", 24, RGBColor.from_string("12395E"),
            True, False, "Definitely Not An Installed Font")]])
    dk.text(s, 0.6, 4.8, 8, 0.3, [[("as of 28 July 2026", 11,
            RGBColor.from_string("646F7B"), False, False, "Definitely Not An Installed Font")]])
    dk.speaker_notes(s, "x")
    prs.save(str(td / "font.pptx"))
    code, out = run(td / "font.pptx", "--static")
    return code == 0 and "not installed on THIS machine" in out and "mechanical subset clean" in out


@case("the suite's own fixtures do not depend on a platform-specific font")
def _(td):
    # A test that passes on macOS and fails on Linux is worse than no test.
    code, out = run(build(td / "plat.pptx"), "--static")
    return code == 0


@case("advisory items are visually distinct from passes")
def _(td):
    code, out = run(build(td / "adv.pptx", asof=False), "--static")
    return "!  7." in out and "✓  8." in out


@case("a clean deck passes")
def _(td):
    code, out = run(build(td / "clean.pptx"), "--static")
    return code == 0 and "mechanical subset clean" in out


@case("a build script that promises build: without Build.step is caught")
def _(td):
    b = td / "build_lying.py"
    b.write_text('def slide01(prs):\n    """build: title then rule"""\n    pass\n')
    code, out = run(build(td / "d1.pptx"), "--static", "--build", str(b))
    return code == 1 and "NO Build.step" in out


@case("an honest build script passes that check")
def _(td):
    b = td / "build_ok.py"
    b.write_text('def slide01(prs):\n    """build: a then b"""\n    with anim.Build(s).step():\n        pass\n')
    code, out = run(build(td / "d2.pptx"), "--static", "--build", str(b))
    return code == 0


@case("an unreadable deck says NOT CHECKED and exits 2, never clean")
def _(td):
    code, out = run(td / "does-not-exist.pptx")
    return code == 2 and "NOT CHECKED" in out and "clean" not in out.lower()


@case("judgment items 5/6/6b/9/11 are printed as NOT covered")
def _(td):
    code, out = run(build(td / "j.pptx"), "--static")
    return all(k in out for k in ("NOT CHECKABLE BY ANY PROGRAM",
                                  "Evidence real", "Eye path", "Titles bound to takeaways"))


@case("the clean run never claims the judgment items were checked")
def _(td):
    code, out = run(build(td / "k.pptx"), "--static")
    tail = out.split("mechanical subset clean")[-1]
    return "still unticked" in tail


def _pitch(td, name, body):
    p = td / name
    p.write_text("import deckkit as dk\n" + body)
    from preflight_check import check_handrolled_pitch
    return check_handrolled_pitch(str(p))


# A literal stride constant in a placement loop is the documented #1 geometry defect, and on one
# measured deck the author called NONE of the derivers and shipped ~15 defects from that
# arithmetic. Both spellings must be caught, and the correct derived form must never be.
@case("a literal `y += 0.72` in a placement loop is caught")
def _(td):
    st, d = _pitch(td, "p1.py",
                   "def s(prs):\n    y = 1.7\n    for a in X:\n"
                   "        dk.box(prs, 0.7, y, 2.5, 0.6)\n        y += 0.72\n")
    return st == "FAIL" and "y += 0.72" in d


@case("`y = base + i * 0.5` is caught through enumerate's TUPLE target")
def _(td):
    st, d = _pitch(td, "p2.py",
                   "def s(prs):\n    for i, (a, b) in enumerate(X):\n"
                   "        yy = 1.68 + i * 0.5\n        dk.icon(prs, 0.7, yy, 0.3)\n")
    return st == "FAIL" and "i * 0.5" in d


@case("a stride DERIVED from the region is never flagged")
def _(td):
    st, _d = _pitch(td, "p3.py",
                    "def s(prs):\n    top, gap, n = 2.9, 0.14, 3\n"
                    "    bh = ((4.1 - top) - gap * (n - 1)) / n\n"
                    "    for i, it in enumerate(X):\n        yy = top + i * (bh + gap)\n"
                    "        dk.box(prs, 5.2, yy, 4.0, bh)\n")
    return st == "PASS"


@case("a function that uses a deriver is left alone")
def _(td):
    st, d = _pitch(td, "p4.py",
                   "def s(prs):\n    for r in dk.rows(3, slide=prs):\n"
                   "        dk.box(prs, *r)\n        y = 1.0 + 0.5\n")
    return st == "PASS" and "rows" in d


def build_grouped(path: Path):
    """A deck carrying two hard-FAIL leaks, with every shape inside a GROUP."""
    import deckkit as dk
    from pptx.dml.color import RGBColor
    C = lambda h: RGBColor.from_string(h)
    prs = dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    seen = set(id(x) for x in s.shapes)
    dk.text(s, 0.6, 1.0, 8, 0.6, [[("Results (editable native chart)", 24, C("12395E"), True, False)]])
    dk.text(s, 0.6, 2.0, 8, 0.6, [[("Cost was <amount> per run", 20, C("12395E"), False, False)]])
    dk.text(s, 0.6, 4.8, 8, 0.3, [[("Figures as of 28 July 2026", 11, C("646F7B"), False, False)]])
    s.shapes.add_group_shape([x for x in list(s.shapes) if id(x) not in seen])
    dk.speaker_notes(s, "spoken thread")
    prs.save(str(path))
    return path


# Every walk in preflight_check stopped at slide.shapes, so a deck whose content lives in GROUPS —
# every designer-tool export, and every deck handed over for redesign — escaped the whole mechanical
# subset and exited 0. Measured before the fix: the top-level walk saw 0 text shapes on this deck.
@case("a meta-annotation leak INSIDE a group is caught")
def _(td):
    code, out = run(build_grouped(td / "grp1.pptx"), "--static")
    return code == 1 and "editable native chart" in out


@case("an unfilled <slot> INSIDE a group is caught")
def _(td):
    code, out = run(build_grouped(td / "grp2.pptx"), "--static")
    return code == 1 and "<amount>" in out


@case("an as-of date INSIDE a group counts as present")
def _(td):
    # the flip side: recursion must not only find faults, it must find the things that CLEAR a check
    code, out = run(build_grouped(td / "grp3.pptx"), "--static")
    return "!  7." not in out


def _means(td, name, body):
    p = td / name
    p.write_text(body)
    from preflight_check import check_bar_of_means
    return check_bar_of_means(str(p))


# A bar's length claims "this value fills the range from zero", which is a statement about a COUNT.
# Fed a mean of measurements it hides n, the spread and the outliers — the one chart choice the
# literature calls a defect rather than a preference. The check keys on STRUCTURE (a computed mean
# reaching a bar), not on variable names, and it must stay silent on the two legitimate neighbours
# or it becomes noise: a bar of counts, and a LINE of averages over time.
@case("a bar/column chart fed a computed mean is flagged")
def _(td):
    st, d = _means(td, "m1.py",
                   "import statistics, deckkit as dk\n"
                   "per = {'a': [.8, .84]}\n"
                   "means = [statistics.mean(v) for v in per.values()]\n"
                   "dk.native_chart(s, 1, 1, 8, 3, list(per), [('Dice', means)], kind='column')\n")
    return st == "ADVISORY" and "distribution" in d


@case("a bar of COUNTS is never flagged")
def _(td):
    st, _d = _means(td, "m2.py",
                    "import deckkit as dk\n"
                    "counts = [12, 40, 7]\n"
                    "dk.native_chart(s, 1, 1, 8, 3, ['a','b','c'], [('n', counts)], kind='column')\n")
    return st == "PASS"


@case("a LINE of averages over time is never flagged")
def _(td):
    st, _d = _means(td, "m3.py",
                    "import numpy as np, deckkit as dk\n"
                    "means = [np.mean(x) for x in series]\n"
                    "dk.native_chart(s, 1, 1, 8, 3, qs, [('avg', means)], kind='line')\n")
    return st == "PASS"


@case("the flag is ADVISORY, not a FAIL (population means are a fair bar)")
def _(td):
    # Nothing in the file separates a sample mean from a population mean, and failing on the
    # latter would train people to ignore the tool — the same reasoning as the as-of-date item.
    st, _d = _means(td, "m4.py",
                    "import statistics, deckkit as dk\n"
                    "avg = statistics.mean(xs)\n"
                    "dk.native_chart(s, 1, 1, 8, 3, c, [('m', avg)], kind='bar')\n")
    return st == "ADVISORY"


@case("no build script is NOT CHECKABLE for the bar-of-means item too")
def _(td):
    from preflight_check import check_bar_of_means
    st, _d = check_bar_of_means(None)
    return st == "NOT CHECKABLE"


@case("a missing build script is NOT CHECKABLE, never clean")
def _(td):
    from preflight_check import check_handrolled_pitch
    st, _d = check_handrolled_pitch(str(td / "nope.py"))
    st2, _d2 = check_handrolled_pitch(None)
    return st == "NOT CHECKABLE" and st2 == "NOT CHECKABLE"


def main() -> int:
    passed = failed = 0
    with tempfile.TemporaryDirectory() as d:
        td = Path(d)
        for name, fn in CASES:
            try:
                ok = fn(td)
            except Exception as e:
                ok = False
                name = f"{name}  [raised {type(e).__name__}: {e}]"
            if ok:
                passed += 1
                print(f"  ok   {name}")
            else:
                failed += 1
                print(f"  FAIL {name}")
    print(f"\n{passed} passed, {failed} failed")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main())
