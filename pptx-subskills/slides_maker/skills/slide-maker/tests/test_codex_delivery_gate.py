#!/usr/bin/env python3
"""Regression fixtures for the strict Codex-only presentation delivery gate.

Run with: python3 tests/test_codex_delivery_gate.py
"""

from __future__ import annotations

import hashlib
import importlib.util
import json
import struct
import sys
import tempfile
import zlib
from pathlib import Path


HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def write_png(path: Path, width: int = 640, height: int = 360, *, alpha: bool = False) -> None:
    def chunk(kind: bytes, payload: bytes) -> bytes:
        return (
            struct.pack(">I", len(payload))
            + kind
            + payload
            + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)
        )

    pixel = b"\xff\xff\xff\xff" if alpha else b"\xff\xff\xff"
    pixels = b"\x00" + pixel * width
    raw = pixels * height
    color_type = 6 if alpha else 2
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, color_type, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw, 9))
        + chunk(b"IEND", b"")
    )


def load_gate():
    spec = importlib.util.spec_from_file_location(
        "codex_delivery_gate_test", SCRIPTS / "codex_delivery_gate.py"
    )
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def fresh_component_audit(build: Path, deck: Path) -> dict:
    spec = importlib.util.spec_from_file_location("component_audit_test", SCRIPTS / "component_audit.py")
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module.audit(str(build), str(deck))


def critic(lens: str) -> dict:
    probes = {"memory_sentence": "The remembered takeaway remains precise and defensible."}
    if lens == "design":
        probes["per_slide"] = [
            {
                "slide": 1,
                "first_read": "A clear title and visual focal point.",
                "takeaway_guess": "The workflow makes evidence review mandatory.",
            }
        ]
    return {
        "purpose": "Validate a one-slide presentation fixture for the final delivery gate.",
        "reviewer": {
            "origin": "isolated",
            "identity": f"fixture-{lens}-reviewer",
            "fresh_context": True,
        },
        "coverage": {
            "slides_opened": [1],
            "passes": [f"{lens} lens (full deck)"],
            "stats_block_seen": True,
            "contract_card_seen": True,
        },
        "plan_audit": {"lens_a": {"status": "pass"}, "lens_b": {"status": "pass"}},
        "probes": probes,
        "verdict": "consent",
        "summary": "The final fixture meets the named review lens.",
        "strengths": ["Evidence and delivery artifacts are bound to the final deck."],
        "findings": [],
    }


def write_json(path: Path, value: dict) -> None:
    path.write_text(json.dumps(value, indent=2) + "\n", encoding="utf-8")


def fixture(root: Path) -> tuple[dict, dict, dict, Path]:
    try:
        from pptx import Presentation
    except ImportError:
        print("SKIPPED: this suite needs python-pptx")
        sys.exit(0)

    deck = root / "deck.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(deck)
    source = root / "README.md"
    source.write_text("# Source\n\nThe gate verifies workflow evidence.\n", encoding="utf-8")
    direction = root / "directions.html"
    direction.write_text(
        "<html><body>" + "direction preview " * 40 + "A B C D</body></html>",
        encoding="utf-8",
    )
    proof = root / "slide-1.png"
    icon = root / "feature.png"
    write_png(proof)
    write_png(icon, 512, 512, alpha=True)
    build = root / "build_deck.py"
    build.write_text(
        "import deckkit as dk\n\n"
        "def slide_01(slide):\n"
        "    dk.icon_card(slide, 0, 0, 1, 1, 'Feature')\n\n"
        "dk.lint_layout(prs, strict=True)\n",
        encoding="utf-8",
    )
    visual_manifest = root / "visual-contract.json"
    write_json(
        visual_manifest,
        {
            "schema": "slide-maker-codex-visual-contract/v1",
            "zones": [],
            "icons": [],
        },
    )
    visual_result = root / "visual-contract-final.json"
    write_json(
        visual_result,
        {
            "schema": "slide-maker-codex-visual-contract-result/v1",
            "pptx_sha256": sha256(deck),
            "manifest_sha256": sha256(visual_manifest),
            "passed": True,
            "zones": [],
            "icons": [],
            "errors": [],
        },
    )
    content_review = root / "critic-content.json"
    design_review = root / "critic-design.json"
    write_json(content_review, critic("content"))
    write_json(design_review, critic("design"))
    evidence = {
        "schema": "slide-maker-codex-evidence/v2",
        "runtime": "codex",
        "delivery": "presented",
        "review_effort": "standard",
        "deck": {"pptx": deck.name, "sha256": sha256(deck), "slide_count": 1},
        "interview": {
            "mode": "answered",
            "record": "The user requested a concise evidence-focused deck.",
        },
        "content": {
            "source_mode": "provided",
            "sources": [{"kind": "provided", "path": source.name, "sha256": sha256(source)}],
            "slides": [
                {
                    "slide": 1,
                    "role": "cover",
                    "takeaway": "The gate makes workflow evidence verifiable.",
                    "evidence": ["README.md:1-3"],
                }
            ],
            "claim_ledger": [
                {
                    "claim": "The gate validates final workflow evidence.",
                    "source": "README.md:3",
                    "verified": True,
                }
            ],
            "checkpoint": {
                "mode": "approved",
                "record": "The one-slide narrative and claim were approved.",
            },
        },
        "design": {
            "direction": {
                "branch": "clean",
                "artifact": direction.name,
                "sha256": sha256(direction),
                "directions": [
                    {
                        "id": "A",
                        "name": "A",
                        "bg": "#071820",
                        "accent": "#54D9D0",
                        "font_display": "Inter",
                        "font_body": "Inter",
                        "density": "minimal",
                        "cover": "low-left",
                        "skeleton": "rail",
                    },
                    {
                        "id": "B",
                        "name": "B",
                        "bg": "#FFFFFF",
                        "accent": "#E4572E",
                        "font_display": "Georgia",
                        "font_body": "Arial",
                        "density": "normal",
                        "cover": "centred",
                        "skeleton": "split",
                    },
                    {
                        "id": "C",
                        "name": "C",
                        "bg": "#182033",
                        "accent": "#F2C14E",
                        "font_display": "Menlo",
                        "font_body": "Inter",
                        "density": "dense",
                        "cover": "full-bleed-type",
                        "skeleton": "dashboard",
                    },
                    {
                        "id": "D",
                        "name": "D",
                        "bg": "#F4E9D8",
                        "accent": "#355C7D",
                        "font_display": "Rockwell",
                        "font_body": "Arial",
                        "density": "spacious",
                        "cover": "split-vertical",
                        "skeleton": "gallery",
                    },
                ],
                "decision": "user-approved",
                "record": "Direction A was selected after reviewing four distinct previews.",
            },
            "type_scale": {"display": 34, "title": 24, "body": 14},
            "boldness": "balanced+",
            "signature_move": "A visible evidence rail connects source, build, and review.",
            "carried_by": [1],
            "signature_proof": {
                "slide": 1,
                "path": proof.name,
                "sha256": sha256(proof),
                "pptx_sha256": sha256(deck),
            },
            "slides": [
                {
                    "slide": 1,
                    "function": "slide_01",
                    "form": "cover",
                    "runner_up": "editorial opener",
                    "reason": "The cover provides a single memorable entry point for the workflow.",
                    "categorical": True,
                    "components": [],
                }
            ],
            "checkpoint": {
                "mode": "approved",
                "record": "The design direction and form ledger were approved.",
            },
        },
        "build": {"script": build.name, "sha256": sha256(build), "strict_layout": True},
        "icons": [
            {
                "slide": 1,
                "family": "lucide",
                "asset": icon.name,
                "sha256": sha256(icon),
                "rasterizer": "scripts/icons.py",
            }
        ],
        "visual_contract": {
            "manifest": visual_manifest.name,
            "sha256": sha256(visual_manifest),
            "result": visual_result.name,
            "result_sha256": sha256(visual_result),
            "pptx_sha256": sha256(deck),
        },
        "critics": [
            {
                "lens": "content",
                "review": content_review.name,
                "sha256": sha256(content_review),
                "pptx_sha256": sha256(deck),
            },
            {
                "lens": "design",
                "review": design_review.name,
                "sha256": sha256(design_review),
                "pptx_sha256": sha256(deck),
            },
        ],
        "waivers": [],
    }
    lint = {
        "findings": [],
        "pixel_checks": [{"pass": True}],
        "text_runs": [{"role": "body", "size_pt": 14}],
        "stats": {"warnings": []},
    }
    components = fresh_component_audit(build, deck)
    return evidence, lint, components, build


def main() -> int:
    gate = load_gate()
    with tempfile.TemporaryDirectory(prefix="codex-gate-") as name:
        root = Path(name)
        evidence, lint, components, build = fixture(root)
        icon = root / evidence["icons"][0]["asset"]
        errors = gate.evaluate(lint, components, build, evidence, root)
        failures = []
        if errors:
            failures.append("valid evidence unexpectedly blocked:\n" + "\n".join(errors))

        evidence["runtime"] = "openai-gpt-bridged"
        errors = gate.evaluate(lint, components, build, evidence, root)
        if errors:
            failures.append("a bridged GPT Store runtime was incorrectly blocked:\n" + "\n".join(errors))
        evidence["runtime"] = "codex"

        lint["stats"]["warnings"] = ["card_dominance"]
        evidence["waivers"] = [
            {
                "kind": "stats",
                "warning": "card_dominance",
                "reason": "The deliberate scorecard composition is the approved signature move.",
            }
        ]
        errors = gate.evaluate(lint, components, build, evidence, root)
        if errors:
            failures.append("a documented, deliberate statistics waiver was blocked:\n" + "\n".join(errors))
        lint["stats"]["warnings"] = []
        evidence["waivers"] = []

        evidence["critics"][0]["pptx_sha256"] = "0" * 64
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("not bound to the final PPTX" in error for error in errors):
            failures.append("a critic review for a stale deck passed the gate")
        evidence["critics"][0]["pptx_sha256"] = evidence["deck"]["sha256"]

        bad_review = root / "critic-content.json"
        write_json(
            bad_review,
            {"purpose": "This artifact is intentionally incomplete but not empty.", "verdict": "consent"},
        )
        evidence["critics"][0]["sha256"] = sha256(bad_review)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("invalid critic schema" in error for error in errors):
            failures.append("an incomplete critic review passed the gate")

        write_json(bad_review, critic("content"))
        evidence["critics"][0]["sha256"] = sha256(bad_review)
        components = {
            "clusters": [{"slide": 1, "pattern": "equal tiles", "suggest": ["scorecard"]}],
            "actionable": [],
            "suppressed_by": ["scorecard"],
        }
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("does not match a fresh audit" in error for error in errors):
            failures.append("a stale or fabricated component audit passed the gate")

        component_errors = []
        gate.check_components(
            evidence,
            {
                "clusters": [{"slide": 1, "pattern": "tile row", "suggest": ["scorecard"]}],
                "suppressed_by": ["org_tree"],
            },
            {1: {"function": "slide_01", "components": ["org_tree"]}},
            {"slide_01": {"org_tree"}},
            component_errors,
        )
        if component_errors:
            failures.append("a same-slide, registered component emitter was incorrectly blocked")

        components = fresh_component_audit(build, root / evidence["deck"]["pptx"])
        proof = root / evidence["design"]["signature_proof"]["path"]
        proof.write_bytes(b"")
        evidence["design"]["signature_proof"]["sha256"] = sha256(proof)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("signature_proof" in error for error in errors):
            failures.append("an empty signature proof passed the gate")

        write_png(icon, 240, 240, alpha=True)
        evidence["icons"][0]["sha256"] = sha256(icon)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("thumbnail blur" in error for error in errors):
            failures.append("a low-resolution icon thumbnail passed the gate")

        write_png(icon, 512, 512, alpha=False)
        evidence["icons"][0]["sha256"] = sha256(icon)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("transparent alpha" in error for error in errors):
            failures.append("a matted icon asset without alpha passed the gate")

        write_png(icon, 512, 512, alpha=True)
        evidence["icons"][0]["sha256"] = sha256(icon)
        build.write_text(
            build.read_text(encoding="utf-8") + "\nimport subprocess\nsubprocess.run(['qlmanage', '-t'], check=False)\n",
            encoding="utf-8",
        )
        evidence["build"]["sha256"] = sha256(build)
        errors = gate.evaluate(lint, components, build, evidence, root)
        if not any("Quick Look thumbnail generation" in error for error in errors):
            failures.append("a Quick Look icon rasterization workaround passed the gate")

        if failures:
            print("\n".join("FAIL: " + failure for failure in failures))
            return 1
    print("ok - Codex delivery gate rejects icon-rasterization and historical bypasses")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
