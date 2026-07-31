#!/usr/bin/env python3
"""Regression checks for Codex-only local visual contracts."""

from __future__ import annotations

import importlib.util
import json
import tempfile
from pathlib import Path


HERE = Path(__file__).resolve().parent
SCRIPT = HERE.parent / "scripts" / "codex_visual_contract.py"


def load_module():
    spec = importlib.util.spec_from_file_location("codex_visual_contract_test", SCRIPT)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def write_json(path: Path, value: dict) -> None:
    path.write_text(json.dumps(value, indent=2) + "\n", encoding="utf-8")


def make_deck(path: Path, *, overlap: bool) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(6.0), Inches(2.0), Inches(2.0), Inches(0.4))
    first.text_frame.text = "交付失控"
    x = 7.2 if overlap else 8.3
    second = slide.shapes.add_textbox(Inches(x), Inches(2.02), Inches(1.4), Inches(0.4))
    second.text_frame.text = "假完成"
    presentation.save(path)


def manifest() -> dict:
    return {
        "schema": "slide-maker-codex-visual-contract/v1",
        "zones": [
            {
                "id": "delivery-callout",
                "slide": 1,
                "subject": {"text": "假完成"},
                "clear_of": [{"text": "交付失控"}],
                "min_gap": 0.12,
            }
        ],
        "icons": [
            {
                "id": "approval-icon",
                "slide": 1,
                "build_token": "lucide:shield-check",
                "meaning": "表示证据已经通过验证并获准交付。",
            }
        ],
    }


def main() -> int:
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("SKIPPED: this suite needs python-pptx")
        return 0
    module = load_module()
    with tempfile.TemporaryDirectory(prefix="codex-visual-contract-") as name:
        root = Path(name)
        build = root / "build.py"
        build.write_text("ICON = 'lucide:shield-check'\n", encoding="utf-8")
        contract = root / "contract.json"
        write_json(contract, manifest())
        bad = root / "bad.pptx"
        make_deck(bad, overlap=True)
        result = module.evaluate(bad, contract, build)
        failures = []
        if result.get("passed") or not any("delivery-callout" in error for error in result.get("errors", [])):
            failures.append("an overlapping titled callout passed the visual contract")

        good = root / "good.pptx"
        make_deck(good, overlap=False)
        result = module.evaluate(good, contract, build)
        if not result.get("passed"):
            failures.append("a clear local callout and semantically bound icon were blocked")

        broken = manifest()
        broken["icons"][0]["build_token"] = "lucide:scan-search"
        write_json(contract, broken)
        result = module.evaluate(good, contract, build)
        if result.get("passed") or not any("approval-icon" in error for error in result.get("errors", [])):
            failures.append("a mismatched icon token passed the visual contract")

        if failures:
            print("\n".join("FAIL: " + failure for failure in failures))
            return 1
    print("ok - Codex visual contracts catch local overlap and icon-semantic drift")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
