#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Two-sided regression test for the render-time lint.

Every check here exists because a REAL deck shipped the defect, or because a real deck was
falsely flagged for craft. The two directions matter equally and are asserted separately:

  PASS deck — ordinary, correctly-built slides plus two DECLARED exceptions (a rhymed triptych,
              a quiet pause page). Zero hard findings. A change that breaks one of these is
              catching craft rather than defects, which is how a rule set makes decks worse.
  FAIL deck — one slide per defect the gates must catch. Each was clean before these checks
              existed; that was the bug.

Run:  python3 tests/test_lint_regressions.py
"""
import os, pathlib, shutil, subprocess, sys, tempfile

HERE = pathlib.Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"


def lint(pptx, renders):
    r = subprocess.run([sys.executable, str(SCRIPTS / "lint_deck.py"), str(pptx),
                        "--renders", str(renders), "--static"],
                       capture_output=True, text=True)
    return r.stdout + r.stderr


def _require_deps():
    """A missing dependency must read as a missing dependency.

    This suite imports deckkit and RENDERS, so it needs python-pptx, Pillow, matplotlib and
    LibreOffice. It was originally wired into CI ABOVE `pip install -r requirements.txt` and
    every run since died on a bare ModuleNotFoundError traceback inside a subprocess — which
    looks exactly like a lint regression and is not one. Fail with a sentence instead.
    """
    missing = []
    for mod, why in (("pptx", "python-pptx"), ("PIL", "Pillow"), ("matplotlib", "matplotlib")):
        try:
            __import__(mod)
        except ImportError:
            missing.append(why)
    if missing:
        print("SKIPPED: this suite needs %s — install requirements.txt BEFORE running it "
              "(in CI, this step must come after the dependency install)." % ", ".join(missing))
        sys.exit(0)
    if not shutil.which("soffice") and not os.path.exists(
            "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        print("SKIPPED: LibreOffice (soffice) not found — this suite renders decks.")
        sys.exit(0)


def main():
    _require_deps()
    tmp = pathlib.Path(tempfile.mkdtemp(prefix="lintfx-"))
    subprocess.run([sys.executable, str(HERE / "lint_fixture.py")], cwd=tmp, check=True,
                   capture_output=True)
    for d in ("fx_pass", "fx_fail", "fx_align_pass", "fx_align_fail"):
        subprocess.run([sys.executable, str(SCRIPTS / "render_deck.py"),
                        str(tmp / f"{d}.pptx"), str(tmp / f"{d}_render")],
                       cwd=tmp, capture_output=True)

    ok, bad = [], []

    def ran(out, label):
        """A token's ABSENCE only means 'suppressed' if the lint actually ran. Without this the
        harness reads a crash as two passing waivers — the same read-silence-as-success mistake
        the checks below exist to catch."""
        if "layout finding(s)" not in out:
            bad.append(f"{label}: lint did not complete, so no assertion below is meaningful:\n"
                       + out.strip()[:400])
            return False
        return True

    # A crashed check must never read as a passed check. The per-slide statistics used to be
    # wrapped in `except Exception: pass`, so one refactor that deleted a local variable took
    # TEXT WALL, LAYOUT SAMENESS, UNDERFILLED and FLAT RHYTHM off every deck while the report
    # still printed "✓ clean" — caught here only by luck, because two of those had tokens
    # asserted below. This asserts the machinery itself, on every deck the suite touches.
    p_out = lint(tmp / "fx_pass.pptx", tmp / "fx_pass_render")
    if not ran(p_out, "PASS deck"):
        print("\n".join("  FAIL " + b for b in bad)); return 1
    if "0 layout finding(s)" in p_out:
        ok.append("PASS deck has zero hard findings")
    else:
        bad.append("PASS deck gained a hard finding — a change is flagging craft, not a defect:\n"
                   + "\n".join(l for l in p_out.splitlines() if ": " in l and "[warn]" not in l))
    # the two DECLARED exceptions must actually be honoured
    if "LAYOUT SAMENESS" in p_out:
        bad.append("declared design_intent(rhyme=) did not suppress LAYOUT SAMENESS — a deliberate "
                   "triptych is being flagged as sameness, and the documented waiver is dead again")
    else:
        ok.append("declared rhyme suppresses LAYOUT SAMENESS")
    if "UNDERFILLED" in p_out:
        bad.append("declared design_intent(envelope=) did not suppress UNDERFILLED — the "
                   "deliberately quiet page cannot be built clean")
    else:
        ok.append("declared quiet envelope suppresses UNDERFILLED")

    f_out = lint(tmp / "fx_fail.pptx", tmp / "fx_fail_render")
    if not ran(f_out, "FAIL deck"):
        print("\n".join("  FAIL " + b for b in bad)); return 1
    for token, what in [
        ("under 3:1, the floor for text at ANY size", "text below the absolute contrast floor is caught"),
        ("RULE THROUGH TEXT", "a hairline painted over type is caught"),
        ("OCCLUSION", "a panel painted over a sentence is caught"),
        ("LAYOUT SAMENESS", "UNdeclared sameness is still flagged (the waiver is not a blanket)"),
        ("UNDERFILLED", "an UNdeclared thin page is still flagged"),
        # The three paint-order faults a PER-SHAPE threshold cannot see. Each shipped on a real
        # deck with the gate reporting clean. Asserted on their own text so a generic finding
        # elsewhere in the deck cannot stand in for them.
        ("OCCLUSION: 'COMPOSITE TILES",
         "150 tiles hiding a caption are caught (union coverage, not one shape at a time)"),
        ("RULE THROUGH TEXT: a", "a rule assembled from 40 dashes is caught like a solid one"),
        ("TEXT NOT VISIBLE: 'an opaque picture over l",
         "a PICTURE over one line is caught from the pixels — unknowable from the XML"),
        # A label must sit on the thing it labels. The panels of a composite figure have no
        # shape geometry, so nothing geometric can see this — and the PASS deck carries the
        # same figures captioned correctly, so over-strictness fails there instead.
        # BOTH slides are asserted by number, and that is the point: the ink test has two
        # halves, slide 11 exercises only `var` and slide 12 only `differs from ground`. With a
        # bare "CAPTION NOT ALIGNED" token, deleting the ground half left the suite green while
        # the real defect (flat panels — an MRI/photo strip) went undetected. Verified by
        # mutation: dropping either half now takes one of these two lines away.
    ]:
        (ok if token in f_out else bad).append(
            what if token in f_out else f"FAIL deck: {token} was NOT raised — the check regressed")

    # ── ALIGNMENT pair, in its own decks. A label must sit on the thing it labels; the panels
    # of a composite figure have no shape geometry, so nothing geometric can see this.
    # Both slides are asserted BY NUMBER because the ink test has two halves and each slide
    # exercises exactly one: slide 1's panels vary down their columns, slide 2's are flat and can
    # only be found by differing from the figure's own ground. With a bare "CAPTION NOT ALIGNED"
    # token, deleting the ground half left the suite green while the real defect — flat panels,
    # i.e. every photo and MRI strip — went undetected. Verified by mutation: dropping either
    # half now takes one of these two lines away.
    a_out = lint(tmp / "fx_align_fail.pptx", tmp / "fx_align_fail_render")
    if ran(a_out, "ALIGN-fail deck"):
        for token, what in [
                ("slide 1: CAPTION NOT ALIGNED",
                 "captions on the text grid under INSET panels are caught (varying columns)"),
                ("slide 2: CAPTION NOT ALIGNED",
                 "the same defect on FLAT panels on their own paper is caught (crop ground)")]:
            (ok if token in a_out else bad).append(
                what if token in a_out
                else f"ALIGN deck: {token} was NOT raised — the check regressed")
    ap_out = lint(tmp / "fx_align_pass.pptx", tmp / "fx_align_pass_render")
    if ran(ap_out, "ALIGN-pass deck"):
        if "CAPTION NOT ALIGNED" in ap_out:
            bad.append("ALIGN-pass deck: correctly centred captions were flagged — the check has "
                       "drifted into over-strictness:\n"
                       + "\n".join(l for l in ap_out.splitlines() if "CAPTION" in l))
        else:
            ok.append("correctly centred captions on both panel shapes stay clean")

    for label, out in (("PASS", p_out), ("FAIL", f_out), ("ALIGN-fail", a_out),
                       ("ALIGN-pass", ap_out)):
        if "[BROKEN]" in out:
            bad.append(f"{label} deck: the per-slide statistics CRASHED — checks silently "
                       f"disabled:\n" + "\n".join(l for l in out.splitlines() if "[BROKEN]" in l))
    if not any("[BROKEN]" in o for o in (p_out, f_out, a_out, ap_out)):
        ok.append("per-slide statistics ran on every fixture deck (no silently-disabled checks)")

    # ── the corpus that matters most: the skill's OWN reference deck. A change that adds a hard
    # finding here is a change that would fail the file SKILL.md tells every builder to copy.
    # This is not hypothetical: promoting the 3.0-4.5 contrast band to a hard failure looked
    # obviously right, passed both synthetic fixtures, and hard-failed this deck four times on
    # accent labels at 4.27:1. Synthetic fixtures cannot tell strictness from correctness; a real
    # deck built by the skill's own example can.
    ex = HERE.parent / "references" / "examples" / "build_example_generic.py"
    if ex.is_file():
        import os
        env = dict(os.environ, PYTHONPATH=str(SCRIPTS))
        r = subprocess.run([sys.executable, str(ex)], cwd=tmp, capture_output=True, text=True, env=env)
        demo = None
        for line in (r.stdout + r.stderr).splitlines():
            if "saved ->" in line:
                demo = pathlib.Path(line.split("saved ->")[1].split("|")[0].strip())
        if demo and demo.is_file():
            subprocess.run([sys.executable, str(SCRIPTS / "render_deck.py"), str(demo),
                            str(tmp / "ex_render")], cwd=tmp, capture_output=True, env=env)
            e_out = subprocess.run([sys.executable, str(SCRIPTS / "lint_deck.py"), str(demo),
                                    "--renders", str(tmp / "ex_render")],
                                   capture_output=True, text=True, env=env).stdout
            n = next((int(l.split(":")[-1].split("layout")[0].strip())
                      for l in e_out.splitlines() if "layout finding(s)" in l), None)
            BASE = 3          # the example deck's own pre-existing footer overlap on slide 4
            if n is None:
                bad.append("reference example deck: lint did not complete")
            elif n <= BASE:
                ok.append(f"reference example deck holds at {n} hard findings (baseline {BASE})")
            else:
                bad.append(f"reference example deck rose to {n} hard findings (baseline {BASE}) — a "
                           f"change is failing the deck the skill tells builders to copy:\n"
                           + "\n".join(l for l in e_out.splitlines()
                                        if ": " in l and "[warn]" not in l and "[stats]" not in l))
        else:
            bad.append("reference example deck could not be built — the corpus check did not run")

    # ── UNSOURCED NUMBER: two-sided, on its own mini-deck ─────────────────────
    # Built separately on purpose: several assertions above name fx_fail slides BY NUMBER, so
    # appending a slide there would silently re-point them. A check that never fires is worth
    # nothing, and one that fires on a sourced deck is worse than nothing, so both directions
    # are asserted. The recap case is the one that matters most — it is the whole reason this
    # check is deck-level: a closing slide restating a figure sourced on its own page is good
    # practice, and a per-slide test calls it a defect.
    sys.path.insert(0, str(SCRIPTS))
    import deckkit as _dk

    def _prov_deck(dest, *slides):
        prs = _dk.blank_deck(10, 5.625)
        for lines, notes in slides:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            for i, ln in enumerate(lines):
                _dk.text(s, 0.6, 0.8 + i * 0.6, 8.6, 0.5,
                         [[(ln, 18, _dk.DEEP, False, False)]])
            if notes:
                _dk.speaker_notes(s, notes)
        prs.save(str(dest))
        return dest

    def _lint_nr(pptx):
        r = subprocess.run([sys.executable, str(SCRIPTS / "lint_deck.py"), str(pptx), "--static"],
                           capture_output=True, text=True)
        return r.stdout + r.stderr

    o = _lint_nr(_prov_deck(tmp / "prov_bare.pptx", (["Capex reached $400B this year"], None)))
    if "UNSOURCED NUMBER" in o and "$400B" in o:
        ok.append("a novel magnitude with no source anywhere is caught")
    else:
        bad.append("UNSOURCED NUMBER did not fire on a bare unsourced $400B — the check is dead")

    o = _lint_nr(_prov_deck(tmp / "prov_notes.pptx",
                            (["Capex reached $400B this year"], "Source: Crunchbase Q1 2026")))
    if "UNSOURCED NUMBER" in o:
        bad.append("UNSOURCED NUMBER fired although the citation is in the SPEAKER NOTES — a "
                   "presented deck legitimately sources there and must not be flagged")
    else:
        ok.append("a citation in the speaker notes counts as provenance")

    o = _lint_nr(_prov_deck(tmp / "prov_recap.pptx",
                            (["Capex reached $400B", "来源: Crunchbase"], None),
                            (["Takeaway: $400B of capex is the story"], None)))
    if "UNSOURCED NUMBER" in o:
        bad.append("UNSOURCED NUMBER fired on a RECAP slide restating a figure sourced earlier in "
                   "the same deck — this is the false positive the deck-level design exists to "
                   "prevent, and it is back")
    else:
        ok.append("a recap of a figure sourced elsewhere in the deck is not flagged")

    o = _lint_nr(_prov_deck(tmp / "prov_chrome.pptx", (["Part III", "13 / 20", "Section 4"], None)))
    if "UNSOURCED NUMBER" in o:
        bad.append("UNSOURCED NUMBER fired on page chrome ('13 / 20') — bare integers are not "
                   "claims, and counting them is what made the first cut fire on a good deck")
    else:
        ok.append("page numbers and section indices are not treated as claims")

    # ── grouped decks must be SEEN, not silently skipped ──────────────────────
    # _boxes() walked slide.shapes only, so on a deck whose content lives in groups — every
    # designer-tool export, and every deck handed over for redesign — it saw one shape per slide.
    # OVERLAP vanished, OVERFLOW could only name "GROUP ''", and every stat derived from boxes
    # (size clusters, text coverage, halves, skeleton, occupancy) read empty while the report still
    # printed "✓ clean". Measured before the fix: identical content, 11 shapes seen ungrouped vs 1
    # grouped.
    def _grp(dest, build, group=True):
        prs = _dk.blank_deck(10, 5.625)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        build(s)
        if group:
            s.shapes.add_group_shape([x for x in list(s.shapes)])
        prs.save(str(dest))
        return dest

    def _defects(s):
        _dk.text(s, 0.6, 0.4, 8.8, 0.6, [[("Title", 28, _dk.DEEP, True, False)]])
        _dk.box(s, 1.0, 1.2, 3.5, 2.2, fill="C0362C")
        _dk.text(s, 7.6, 4.6, 3.2, 0.5, [[("runs off the right edge", 13, _dk.DEEP, False, False)]])

    import sys as _sys
    _sys.path.insert(0, str(SCRIPTS))
    import lint_deck as _L
    from pptx import Presentation as _P

    a = sorted((round(b["l"], 3), round(b["t"], 3), round(b["w"], 3), round(b["h"], 3))
               for b in _L._boxes(_P(str(_grp(tmp / "g_flat.pptx", _defects, False))).slides[0], 10, 5.625))
    b = sorted((round(x["l"], 3), round(x["t"], 3), round(x["w"], 3), round(x["h"], 3))
               for x in _L._boxes(_P(str(_grp(tmp / "g_grp.pptx", _defects, True))).slides[0], 10, 5.625))
    if a and a == b:
        ok.append("grouped and ungrouped decks yield IDENTICAL geometry (transform applied)")
    else:
        bad.append(f"grouped geometry does not match ungrouped: {len(a)} vs {len(b)} boxes; the "
                   f"group transform (off/ext/chOff/chExt) is wrong or the walk stopped at the group")

    o = _lint_nr(tmp / "g_grp.pptx")
    if "OVERFLOW" in o and "GROUP" not in o.split("OVERFLOW")[1][:40]:
        ok.append("OVERFLOW inside a group names the real shape, not the group")
    else:
        bad.append("OVERFLOW on a grouped deck did not resolve to the offending child shape")

    # A group is an AUTHORED unit: a badge straddling a card corner is layering, not a collision.
    def _composed(s):
        _dk.box(s, 0.8, 1.5, 3.6, 2.2, fill="FFFFFF", line="DDDDDD")
        _dk.box(s, 4.0, 1.32, 1.1, 0.42, fill="C0362C")      # badge over the card's top-right
        _dk.text(s, 4.05, 1.36, 1.0, 0.34, [[("NEW", 10, _dk.WHITE, True, False)]])
    o = _lint_nr(_grp(tmp / "g_composed.pptx", _composed, True))
    if "OVERLAP" in o:
        bad.append("OVERLAP fired INSIDE one group — a group is an authored composition (badge on a "
                   "card, icon on a panel); flagging it makes the check useless on designed decks")
    else:
        ok.append("layering INSIDE one group is not flagged as a collision")

    def _cross(s):
        _dk.box(s, 1.0, 1.6, 3.0, 1.6, fill="F2F4F8")
        _dk.text(s, 1.2, 1.8, 2.6, 0.5, [[("card", 12, _dk.DEEP, False, False)]])
    prs = _dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _cross(s)
    s.shapes.add_group_shape([x for x in list(s.shapes)])
    _dk.box(s, 2.2, 2.4, 3.4, 2.0, fill="C0362C")            # OUTSIDE the group
    prs.save(str(tmp / "g_cross.pptx"))
    if "OVERLAP" in _lint_nr(tmp / "g_cross.pptx"):
        ok.append("a collision ACROSS a group boundary is still caught")
    else:
        bad.append("a shape colliding with a group's contents is no longer caught — the same-group "
                   "exemption has widened past the composed unit it was carved for")

    # An unmappable group must be reported ONCE, and must not take the run down with it. Both of
    # these were live bugs: the stats walks re-traversed the same tree and re-reported every rotated
    # group as "slide ?", and _report_group_skip sorted int slide numbers against that None, raising
    # a TypeError at the very END of lint() — discarding a whole completed run.
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    prs = _dk.blank_deck(10, 5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dk.box(s, 1.0, 1.0, 2.0, 1.0, fill="C0362C")
    g = s.shapes.add_group_shape([x for x in list(s.shapes)])
    g._element.find(f".//{A_NS}xfrm").set("rot", "2700000")
    prs.save(str(tmp / "g_rot.pptx"))
    o = _lint_nr(tmp / "g_rot.pptx")
    n_rot = o.count("rotated/flipped group")
    if n_rot == 1:
        ok.append("an unmappable group is reported exactly once, with its slide number")
    else:
        bad.append(f"the rotated-group skip was printed {n_rot}x (want 1) — the stats walks are "
                   f"re-recording it, which also reintroduces the None-vs-int sort crash")
    if "layout finding(s)" in o and "slide ?" not in o:
        ok.append("lint completes on a deck whose group cannot be mapped")
    else:
        bad.append("lint did not complete cleanly on an unmappable-group deck (or printed 'slide ?')")

    # ── the declared type_scale must actually bind ────────────────────────────
    # render_deck --gate-check requires design_plan.type_scale; nothing compared it to the deck, so
    # a deck could declare {34,24,14}, set 31/22/17 throughout, and pass both gates clean (measured).
    # The rule has to stay narrow: the skill's own 5-slide example uses TWELVE distinct sizes, so
    # "every size must be a declared tier" would fire on correct work and be abandoned immediately.
    import json as _json

    def _scaled(dest, sizes, declared):
        prs = _dk.blank_deck(10, 5.625)
        for i in range(3):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _dk.text(s, 0.6, 0.5, 8.8, 0.7, [[(f"Title {i+1}", sizes[0], _dk.DEEP, True, False)]])
            _dk.text(s, 0.6, 1.5, 8.8, 0.5, [[("Subhead", sizes[1], _dk.DEEP, False, False)]])
            _dk.text(s, 0.6, 2.3, 8.8, 1.5,
                     [[("Body copy carrying most of this deck's characters by a wide margin.",
                        sizes[2], _dk.MUTE, False, False)]])
        d = dest.parent / dest.stem
        d.mkdir(exist_ok=True)
        out = d / "t.pptx"
        prs.save(str(out))
        if declared:
            (d / ".deck-gates.json").write_text(_json.dumps({"design_plan": {"type_scale": declared}}))
        return out

    o = _lint_nr(_scaled(tmp / "sc_ok.pptx", (34, 24, 14), {"display": 34, "title": 24, "body": 14}))
    if "SCALE DRIFT" in o:
        bad.append("SCALE DRIFT fired on a deck that sets exactly what it declares")
    else:
        ok.append("a deck that honours its declared type scale is silent")

    o = _lint_nr(_scaled(tmp / "sc_drift.pptx", (31, 22, 17), {"display": 34, "title": 24, "body": 14}))
    if "SCALE DRIFT" in o and "body=14" in o:
        ok.append("a deck that declares one scale and sets another is caught")
    else:
        bad.append("SCALE DRIFT missed a deck declaring 34/24/14 while setting 31/22/17 — the "
                   "required field constrains nothing again")

    o = _lint_nr(_scaled(tmp / "sc_none.pptx", (34, 24, 14), None))
    if "SCALE DRIFT" in o:
        bad.append("SCALE DRIFT fired with no .deck-gates.json — it must be silent when nothing "
                   "was declared, or every deck without a gates file gains a warning")
    else:
        ok.append("no declaration means no drift finding")

    # a hero number way off the scale is normal and must not read as drift
    prs = _dk.blank_deck(10, 5.625)
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _dk.text(s, 0.6, 0.5, 8.8, 0.8, [[(f"Title {i+1}", 34, _dk.DEEP, True, False)]])
        _dk.text(s, 0.6, 1.6, 8.8, 0.6, [[("Subhead", 24, _dk.DEEP, False, False)]])
        _dk.text(s, 0.6, 2.4, 8.8, 1.6,
                 [[("Body copy carrying most of the characters in this deck.", 14, _dk.MUTE, False, False)]])
    s = prs.slides[1]
    _dk.text(s, 6.6, 3.6, 3.0, 1.2, [[("97%", 96, _dk.DEEP, True, False)]])   # deliberate hero
    d = tmp / "sc_hero"; d.mkdir(exist_ok=True)
    prs.save(str(d / "t.pptx"))
    (d / ".deck-gates.json").write_text(_json.dumps(
        {"design_plan": {"type_scale": {"display": 34, "title": 24, "body": 14}}}))
    if "SCALE DRIFT" in _lint_nr(d / "t.pptx"):
        bad.append("SCALE DRIFT fired on a deliberate off-scale hero number — the check must key on "
                   "the size carrying the TEXT, not on every size present")
    else:
        ok.append("a deliberate off-scale hero number is not drift")

    # A run can INHERIT its size from the layout/theme and report none, so a template-based deck can
    # leave most of its body unmeasurable. Judging "which size carries the most text" from what is
    # left reads a caption as the body: measured, a deck with 1200 inherited characters and one
    # 28-character 9pt caption produced THREE confident findings from those 28 characters.
    prs = _dk.blank_deck(10, 5.625)
    from pptx.util import Inches as _In
    for _i in range(4):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(_In(0.6), _In(1), _In(8.6), _In(2))
        tb.text_frame.text = "Body copy that inherits its size from the theme. " * 6
        _dk.text(s, 0.6, 3.6, 3, 0.4, [[("caption", 9, _dk.MUTE, False, False)]])
    d = tmp / "sc_inherit"; d.mkdir(exist_ok=True)
    prs.save(str(d / "t.pptx"))
    (d / ".deck-gates.json").write_text(_json.dumps(
        {"design_plan": {"type_scale": {"display": 34, "title": 24, "body": 14}}}))
    o = _lint_nr(d / "t.pptx")
    if "SCALE DRIFT NOT CHECKED" in o:
        ok.append("a deck whose type is mostly INHERITED refuses to be judged, and says so")
    elif "SCALE DRIFT" in o:
        bad.append("SCALE DRIFT judged a deck from a thin sample of explicitly-sized text — a "
                   "caption was read as the body, which is a confident wrong finding")
    else:
        bad.append("a deck with unmeasurable type went silent instead of reporting NOT CHECKED — "
                   "silence here is indistinguishable from a pass")

    # ── the APPENDIX run ──────────────────────────────────────────────────────
    # A thesis defense is told to "plan for backup/appendix slides for Q&A", and those are dense ON
    # PURPOSE. Judged as presented content they drew TEXT WALL + CROWDED on every one (measured: 6
    # findings on 3 backup slides), and the trailing run also stole the closing slide's exemption by
    # making a backup slide the last one. Declaring the run must fix both — and must not become a
    # free pass for cramming.
    def _defense(dest, declare):
        prs = _dk.blank_deck(10, 5.625)
        _dk.text(prs.slides.add_slide(prs.slide_layouts[6]), 1, 2, 8, 1,
                 [[("Defense", 34, _dk.DEEP, True, False)]])
        for k in range(2):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _dk.text(s, 0.6, 0.5, 8.8, 0.6, [[(f"Result {k+1}", 24, _dk.DEEP, True, False)]])
            _dk.text(s, 0.6, 1.5, 8.8, 1.0, [[("One clear finding.", 15, _dk.MUTE, False, False)]])
        s = prs.slides.add_slide(prs.slide_layouts[6])          # the real closer: thin by design
        # >=15 words on purpose: UNDERFILLED only looks at slides carrying real text, so a shorter
        # closer would never exercise the exemption this test exists to check.
        _dk.text(s, 0.6, 1.8, 8.8, 1.2,
                 [[("Contributions: a faithful reconstruction method, validated on three separate "
                    "cohorts, with limitations stated plainly and the code released for reproduction.",
                    18, _dk.DEEP, True, False)]])
        for k in range(3):                                       # dense backup material
            b = prs.slides.add_slide(prs.slide_layouts[6])
            if declare and k == 0:
                _dk.design_intent(b, role="appendix", reason="backup slides for Q&A")
            _dk.text(b, 0.6, 0.4, 8.8, 0.5, [[(f"Backup {k+1}", 20, _dk.DEEP, True, False)]])
            _dk.text(b, 0.6, 1.1, 8.8, 4.0,
                     [[(("Full ablation detail with every hyperparameter, dataset split, random seed "
                         "and per-cohort breakdown, kept off the main line but ready. ") * 4,
                        12, _dk.MUTE, False, False)]])
        prs.save(str(dest))
        return dest

    o = _lint_nr(_defense(tmp / "ap_undeclared.pptx", False))
    if "TEXT WALL" in o and "UNDERFILLED" in o:
        ok.append("an UNDECLARED trailing appendix is still judged as presented content")
    else:
        bad.append("the appendix exemption applies without being declared — it must not be free")

    o = _lint_nr(_defense(tmp / "ap_declared.pptx", True))
    if "TEXT WALL" in o:
        bad.append("dense backup slides still draw TEXT WALL after declaring role='appendix' — "
                   "reference material read on demand is dense on purpose")
    elif "UNDERFILLED" in o:
        bad.append("the closing slide still loses its exemption to a trailing appendix — last_body "
                   "is not being honoured")
    else:
        ok.append("a declared appendix reads at briefing density and gives the closer back its exemption")

    # the exemption RAISES the bar, it does not remove it
    prs = _dk.blank_deck(10, 5.625)
    _dk.text(prs.slides.add_slide(prs.slide_layouts[6]), 1, 2, 8, 1,
             [[("Cover", 34, _dk.DEEP, True, False)]])
    for k in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _dk.text(s, 0.6, 0.5, 8.8, 0.6, [[(f"R{k+1}", 24, _dk.DEEP, True, False)]])
        _dk.text(s, 0.6, 1.5, 8.8, 1.0, [[("One finding.", 15, _dk.MUTE, False, False)]])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dk.text(s, 0.6, 1.8, 8.8, 1.0, [[("Contributions and limitations owned.", 20, _dk.DEEP, True, False)]])
    b = prs.slides.add_slide(prs.slide_layouts[6])
    _dk.design_intent(b, role="appendix", reason="Q&A")
    _dk.box(b, 0.2, 0.2, 9.6, 5.2, fill="EEF2F7")
    _dk.text(b, 0.3, 0.3, 9.4, 5.0,
             [[(("Freeform cramming with no structure at all. " * 30), 11, _dk.MUTE, False, False)]])
    prs.save(str(tmp / "ap_crammed.pptx"))
    if "TEXT WALL" in _lint_nr(tmp / "ap_crammed.pptx"):
        ok.append("a CRAMMED appendix slide is still caught — the bar rises, it does not vanish")
    else:
        bad.append("role='appendix' silenced a genuinely crammed slide — dense is correct there, "
                   "freeform cramming is not")

    # ── the translucent-overlap exemption must stay NARROW ────────────────────
    # venn()'s circles overlap by definition, so OVERLAP fired on every Venn ever built — a hard
    # gate failing on correct output is worse than no gate, because the author starts working
    # around the linter. The exemption keys on translucency (a gradient fill is this toolkit's only
    # alpha path) PLUS equal size. All four directions are asserted: widening it later would
    # silently switch off the card-on-card collision check, which is the whole point of OVERLAP.
    def _ov(dest, build):
        prs = _dk.blank_deck(10, 5.625)
        build(prs.slides.add_slide(prs.slide_layouts[6]))
        prs.save(str(dest))
        return dest

    def _g(c):
        return [(0.0, c, 0.3), (1.0, c, 0.3)]

    o = _lint_nr(_ov(tmp / "ov_opaque.pptx", lambda s: (
        _dk.box(s, 1.0, 1.0, 3.0, 2.0, fill="C0362C"),
        _dk.box(s, 2.2, 1.6, 3.0, 2.0, fill="1F77C4"))))
    if "OVERLAP" in o:
        ok.append("an OPAQUE same-size partial overlap still fires")
    else:
        bad.append("OVERLAP stopped firing on two opaque same-size cards — the translucency "
                   "exemption has widened into the defect it was carved around")

    o = _lint_nr(_ov(tmp / "ov_venn.pptx", lambda s: _dk.venn(
        s, 0.6, 1.0, 5.0, 3.8, ["A", "B"], zones={"12": "both"})))
    if "OVERLAP" in o:
        bad.append("OVERLAP fired on venn()'s circles — a Venn's regions ARE overlaps, so the gate "
                   "fails correct output and cannot be satisfied")
    else:
        ok.append("venn()'s intentional circle overlap is exempt")

    o = _lint_nr(_ov(tmp / "ov_mismatch.pptx", lambda s: (
        _dk.box(s, 1.0, 1.0, 4.0, 2.6, grad=_g("C0362C")),
        _dk.box(s, 3.6, 2.2, 2.2, 1.4, grad=_g("1F77C4")))))
    if "OVERLAP" in o:
        ok.append("translucent but MISMATCHED sizes still fires (not a set diagram)")
    else:
        bad.append("OVERLAP stopped firing on translucent shapes of different sizes — equal size is "
                   "what makes the exemption a set diagram rather than a stray panel")

    o = _lint_nr(_ov(tmp / "ov_text.pptx", lambda s: (
        _dk.box(s, 1.0, 1.0, 3.0, 2.0, grad=_g("C0362C")),
        _dk.text(s, 1.2, 1.4, 2.6, 0.6, [[("covered sentence", 14, _dk.DEEP, False, False)]]),
        _dk.box(s, 1.0, 1.0, 3.0, 2.0, grad=_g("1F77C4")))))
    if "layout finding(s)" in o:
        ok.append("a translucent pair carrying TEXT is still assessed (exemption needs textless)")
    else:
        bad.append("the lint did not complete on the translucent-plus-text deck")

    # ── text MEASUREMENT must never under-estimate ────────────────────────────
    # macOS ships a whole family as one .ttc, and matplotlib resolves bold and regular to
    # the SAME path; Pillow's truetype(path, size) with no index= then loads face 0, the
    # Regular. Every bold run in such a family was measured at REGULAR width — 3.9% narrow
    # for Helvetica Neue. Under-measuring is the dangerous direction: a measure-then-place
    # guard silently passes and the renderer wraps anyway, which is how a caption sized for
    # one line landed its second line on top of a footer, on a deck the lint called clean.
    sys.path.insert(0, str(SCRIPTS))
    import deckkit as dk
    fams = [f for f in ("Helvetica Neue", "Arial", "Helvetica", "DejaVu Sans")
            if dk._font_file(f, False)]
    if not fams:
        ok.append("font measurement: skipped — no resolvable font on this host")
    for fam in fams:
        probe = "MEASURING WIDTH AT A GIVEN WEIGHT"
        reg = dk._pil_font(fam, 100, False).getlength(probe)
        bld = dk._pil_font(fam, 100, True).getlength(probe)
        if bld < reg:
            bad.append(f"{fam}: bold measures NARROWER than regular ({bld:.0f} < {reg:.0f}) — "
                       f"the bold face was not selected inside the font file")
        else:
            ok.append(f"{fam}: bold measures >= regular ({bld / reg:.3f}x)")
        path = str(dk._font_file(fam, True) or "")
        if path.lower().endswith((".ttc", ".otc")) and dk._face_index(path, True) == 0:
            bad.append(f"{fam} is a font collection but bold still maps to face 0")

    # ── the measurement must agree with the RENDERER, not just with itself ────
    # Every measure-then-place guard in this skill — head()'s title assert, bound()'s caption
    # sizing, `assert measure_text(...) < h`, and _rbox(), which every geometry check is
    # computed against — trusts one number: how wide the renderer will set this string. When
    # that number came back 3.9% narrow for bold text in font-collection families, all of them
    # silently PASSED while text wrapped anyway, and a caption sized for one line put its
    # second line on top of a footer. The checks above cannot see that: they are computed from
    # the same wrong number. So calibrate against the only authority there is — render real
    # strings and measure the ink. Under-measuring is the dangerous direction; over-measuring
    # only costs a little slack, so the bar is one-sided and tight on the side that hurts.
    from PIL import Image
    cal_fonts = [f for f in ("Helvetica Neue", "Arial", "Helvetica", "DejaVu Sans")
                 if dk._font_file(f, False)]
    CASES = [("Measuring width at a given weight", 10, False),
             ("MEASURING WIDTH AT A GIVEN WEIGHT", 10, True),
             ("Neither number tells you how this compares", 13, False),
             ("Both numbers compare one thing to another.", 26, True),
             ("Ask for the comparator.", 48, True),
             ("1,356", 96, True)]
    for fam in cal_fonts[:2]:                       # two families is enough to catch the class
        prs = dk.blank_deck()
        for txt, size, bold in CASES:
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            dk.box(sl, 0, 0, 10, 5.625, fill=dk.WHITE, line=None)
            dk.text(sl, 0.05, 2.2, 9.9, size / 72.0 * 1.6,
                    [[(txt, size, dk.DEEP, bold, False, fam)]], space_after=0, wrap=False)
        mp = tmp / ("cal_%s.pptx" % fam.replace(" ", ""))
        prs.save(str(mp))
        subprocess.run([sys.executable, str(SCRIPTS / "render_deck.py"), str(mp),
                        str(tmp / ("cal_%s" % fam.replace(" ", "")))],
                       cwd=tmp, capture_output=True)
        worst = None
        for i, (txt, size, bold) in enumerate(CASES, 1):
            png = tmp / ("cal_%s" % fam.replace(" ", "")) / ("slide%02d.png" % i)
            if not png.is_file():
                continue
            im = Image.open(png).convert("L").point(lambda v: 255 if v < 128 else 0)
            bb = im.getbbox()
            if not bb:
                continue
            rendered = (bb[2] - bb[0]) / (im.width / 10.0)
            measured = dk._pil_font(fam, size, bold).getlength(txt) / dk._MEAS_PREC / 72.0
            if measured <= 0:
                continue
            r = rendered / measured
            if worst is None or r > worst[0]:
                worst = (r, txt, size, bold)
        if worst is None:
            bad.append(f"font measurement vs render: {fam} produced no usable renders")
        elif worst[0] > 1.005:
            bad.append(f"{fam}: the renderer sets text {100*(worst[0]-1):.1f}% WIDER than the "
                       f"measurement predicts (worst: {worst[2]}pt "
                       f"{'bold' if worst[3] else 'regular'}, {worst[1]!r}). Every "
                       f"measure-then-place guard in the skill silently passes when this drifts")
        elif worst[0] < 0.80:
            bad.append(f"{fam}: measurement over-estimates width by more than 20% "
                       f"(ratio {worst[0]:.3f}) — layouts will be needlessly cramped")
        else:
            ok.append(f"{fam}: measurement matches the renderer (worst ratio {worst[0]:.3f}, "
                      f"never under-measures)")

    for line in ok:
        print("  ok   " + line)
    for line in bad:
        print("  FAIL " + line)
    print(f"\n{len(ok)} passed, {len(bad)} failed")
    return 1 if bad else 0


if __name__ == "__main__":
    sys.exit(main())
