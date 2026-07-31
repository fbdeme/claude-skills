#!/usr/bin/env python3
"""Strict, Codex-only delivery gate for the slide-maker workflow.

This script intentionally does not change deckkit, component_audit, or the
Claude Code workflow. It verifies the additional evidence Codex must collect
before claiming a deck is ready for delivery.
"""

from __future__ import annotations

import argparse
import ast
import hashlib
import importlib.util
import json
import os
import struct
import subprocess
import sys
from collections import defaultdict
from pathlib import Path
from typing import Any


SCHEMA = "slide-maker-codex-evidence/v2"
BODY_FLOORS = {"presented": 13.5, "textheavy": 13.5, "selfread": 12.0}
STRICT_STATS = {
    "card_dominance",
    "cjk_risk",
    "color_envelope",
    "flat_type",
    "size_sprawl",
    "small_type",
    "timid_cover",
}
ICON_HELPERS = {"icon", "icon_card", "icon_tile", "icon_badge", "icon_ghost"}

TEMPLATE = {
    "schema": SCHEMA,
    "runtime": "codex",
    "delivery": "presented",
    "review_effort": "standard",
    "fast_opt_in": "<only when review_effort is fast: >=12 chars recording the user asking for it>",
    "deck": {
        "pptx": "deck.pptx",
        "sha256": "<sha256 of final deck.pptx>",
        "slide_count": 10,
    },
    "interview": {
        "mode": "answered",
        "record": "<user answers or auto-carved rationale>",
    },
    "content": {
        "source_mode": "provided",
        "sources": [
            {
                "kind": "provided",
                "path": "README.md",
                "sha256": "<sha256>",
            }
        ],
        "slides": [
            {
                "slide": 1,
                "role": "cover",
                "takeaway": "<one audience takeaway>",
                "evidence": ["README.md:1-10"],
            }
        ],
        "claim_ledger": [
            {
                "claim": "<checkable claim>",
                "source": "README.md:1-10",
                "verified": True,
            }
        ],
        "checkpoint": {"mode": "approved", "record": "<decision record>"},
    },
    "design": {
        "direction": {
            "branch": "clean",
            "artifact": "directions.html",
            "sha256": "<sha256>",
            "directions": [
                {
                    "id": "A",
                    "name": "A",
                    "bg": "#071820",
                    "accent": "#54D9D0",
                    "font_display": "Inter",
                    "font_body": "Inter",
                    "density": "minimal",
                    "cover": "low-left",
                    "skeleton": "rail",
                }
            ],
            "decision": "user-approved",
            "record": "<selection or auto-carve record>",
        },
        "type_scale": {"display": 34, "title": 24, "body": 14},
        "boldness": "balanced+",
        "signature_move": "<repeated, deliberate visual device>",
        "carried_by": [1, 5],
        "signature_proof": {
            "slide": 1,
            "path": "render/slide-1.png",
            "sha256": "<sha256>",
            "pptx_sha256": "<must equal deck.sha256>",
        },
        "slides": [
            {
                "slide": 1,
                "function": "slide_01",
                "form": "cover",
                "runner_up": "editorial opener",
                "reason": "<why this form serves the takeaway>",
                "categorical": False,
                "components": [],
            }
        ],
        "checkpoint": {"mode": "approved", "record": "<decision record>"},
    },
    "build": {
        "script": "build_deck.py",
        "sha256": "<sha256>",
        "strict_layout": True,
    },
    "icons": [
        {
            "slide": 2,
            "family": "lucide",
            "asset": "assets/icons/feature.png",
            "sha256": "<sha256>",
            "rasterizer": "scripts/icons.py",
        }
    ],
    "visual_contract": {
        "manifest": "visual-contract.json",
        "sha256": "<sha256>",
        "result": "visual-contract-final.json",
        "result_sha256": "<sha256>",
        "pptx_sha256": "<must equal deck.sha256>",
    },
    "critics": [
        {
            "lens": "content",
            "review": "critic-content-round2.json",
            "sha256": "<sha256>",
            "pptx_sha256": "<must equal deck.sha256>",
        },
        {
            "lens": "design",
            "review": "critic-design-round2.json",
            "sha256": "<sha256>",
            "pptx_sha256": "<must equal deck.sha256>",
        },
    ],
    "thorough_panel": None,
    "arbiters": [],
    "waivers": [],
}


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        value = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return value


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def resolve_path(root: Path, value: Any) -> Path | None:
    if not isinstance(value, str) or not value.strip():
        return None
    path = Path(value)
    return path if path.is_absolute() else root / path


def is_sha256(value: Any) -> bool:
    if not isinstance(value, str) or len(value) != 64:
        return False
    return all(char in "0123456789abcdef" for char in value.lower())


def require_string(value: Any, label: str, errors: list[str], minimum: int = 1) -> str | None:
    if not isinstance(value, str) or len(value.strip()) < minimum:
        errors.append(f"{label} must be a non-empty string")
        return None
    return value.strip()


def check_hashed_file(
    root: Path,
    path_value: Any,
    expected_hash: Any,
    label: str,
    errors: list[str],
    minimum_bytes: int = 1,
) -> Path | None:
    path = resolve_path(root, path_value)
    if path is None:
        errors.append(f"{label}.path must be a non-empty path")
        return None
    if not path.is_file():
        errors.append(f"{label} file not found: {path}")
        return None
    if path.stat().st_size < minimum_bytes:
        errors.append(f"{label} is too small to be usable: {path}")
        return None
    if not is_sha256(expected_hash):
        errors.append(f"{label}.sha256 must be a 64-character SHA-256")
        return path
    actual_hash = sha256_file(path)
    if actual_hash != expected_hash:
        errors.append(f"{label} SHA-256 does not match: {path}")
    return path


def png_dimensions(path: Path) -> tuple[int, int] | None:
    info = png_info(path)
    return info[:2] if info is not None else None


def png_info(path: Path) -> tuple[int, int, int, int] | None:
    """Return PNG width, height, bit depth, and colour type without a Pillow dependency."""
    try:
        with path.open("rb") as handle:
            header = handle.read(26)
    except OSError:
        return None
    if len(header) != 26 or header[:8] != b"\x89PNG\r\n\x1a\n" or header[12:16] != b"IHDR":
        return None
    width, height = struct.unpack(">II", header[16:24])
    return width, height, header[24], header[25]


def slide_count_from_pptx(path: Path) -> int | None:
    try:
        from pptx import Presentation

        return len(Presentation(str(path)).slides)
    except Exception:
        return None


def waived(evidence: dict[str, Any], kind: str, **matches: Any) -> bool:
    waivers = evidence.get("waivers", [])
    if not isinstance(waivers, list):
        return False
    for entry in waivers:
        if not isinstance(entry, dict) or entry.get("kind") != kind:
            continue
        if not require_string(entry.get("reason"), "waiver.reason", [], minimum=12):
            continue
        if all(entry.get(key) == value for key, value in matches.items()):
            return True
    return False


def parse_script(script_path: Path, errors: list[str]) -> tuple[dict[str, set[str]], bool]:
    try:
        tree = ast.parse(script_path.read_text(encoding="utf-8"), filename=str(script_path))
    except (OSError, SyntaxError) as exc:
        errors.append(f"cannot parse build script {script_path}: {exc}")
        return {}, False

    module_aliases: set[str] = {"deckkit"}
    direct_aliases: dict[str, str] = {}
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for item in node.names:
                if item.name == "deckkit":
                    module_aliases.add(item.asname or "deckkit")
        elif isinstance(node, ast.ImportFrom) and node.module == "deckkit":
            for item in node.names:
                direct_aliases[item.asname or item.name] = item.name

    def call_name(node: ast.expr) -> str | None:
        if isinstance(node, ast.Name):
            return direct_aliases.get(node.id, node.id)
        if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name):
            if node.value.id in module_aliases:
                return node.attr
        return None

    calls: dict[str, set[str]] = defaultdict(set)
    strict_layout = False

    class Visitor(ast.NodeVisitor):
        current_function: str | None = None

        def visit_FunctionDef(self, node: ast.FunctionDef) -> None:
            previous = self.current_function
            self.current_function = node.name
            self.generic_visit(node)
            self.current_function = previous

        visit_AsyncFunctionDef = visit_FunctionDef

        def visit_Call(self, node: ast.Call) -> None:
            nonlocal strict_layout
            name = call_name(node.func)
            if self.current_function and name:
                calls[self.current_function].add(name)
            if name == "lint_layout":
                strict_layout = any(
                    keyword.arg == "strict"
                    and isinstance(keyword.value, ast.Constant)
                    and keyword.value.value is True
                    for keyword in node.keywords
                ) or strict_layout
            self.generic_visit(node)

    Visitor().visit(tree)
    return dict(calls), strict_layout


def forbidden_icon_rasterizer_calls(script_path: Path, errors: list[str]) -> list[str]:
    """Find actual command invocations, not comments, that make preview thumbnails into icons."""
    try:
        tree = ast.parse(script_path.read_text(encoding="utf-8"), filename=str(script_path))
    except (OSError, SyntaxError) as exc:
        errors.append(f"cannot inspect build script for icon rasterization: {exc}")
        return []

    command_methods = {"run", "Popen", "call", "check_call", "check_output", "system", "popen"}
    blocked: list[str] = []
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call):
            continue
        if isinstance(node.func, ast.Name):
            command_call = node.func.id in command_methods
        elif isinstance(node.func, ast.Attribute):
            command_call = node.func.attr in command_methods
        else:
            command_call = False
        if not command_call:
            continue
        strings = [
            value.value.lower()
            for value in ast.walk(node)
            if isinstance(value, ast.Constant) and isinstance(value.value, str)
        ]
        if any("qlmanage" in value for value in strings):
            blocked.append("qlmanage")
    return sorted(set(blocked))


def check_lint(lint: dict[str, Any], delivery: str, evidence: dict[str, Any], errors: list[str]) -> None:
    findings = lint.get("findings", [])
    if not isinstance(findings, list):
        errors.append("lint findings missing or malformed")
    else:
        blocking = [finding for finding in findings if isinstance(finding, dict) and finding.get("severity") == "error"]
        if blocking:
            errors.append(f"layout lint reports {len(blocking)} error(s)")

    pixels = lint.get("pixel_checks", [])
    if not isinstance(pixels, list):
        errors.append("pixel_checks missing or malformed")
    else:
        failed = [item for item in pixels if isinstance(item, dict) and not item.get("pass", False)]
        if failed:
            errors.append(f"pixel checks report {len(failed)} failure(s)")

    floor = BODY_FLOORS.get(delivery, BODY_FLOORS["presented"])
    text_runs = lint.get("text_runs", [])
    if not isinstance(text_runs, list):
        errors.append("text_runs missing or malformed")
    else:
        undersized = [
            run
            for run in text_runs
            if isinstance(run, dict)
            and run.get("role") in {"body", "label", "footer"}
            and isinstance(run.get("size_pt"), (int, float))
            and run["size_pt"] < floor
            and not run.get("exception")
        ]
        if undersized:
            errors.append(f"{len(undersized)} body/label/footer run(s) are below {floor:g} pt")

    stats = lint.get("stats", {})
    warnings = stats.get("warnings", []) if isinstance(stats, dict) else []
    if not isinstance(warnings, list):
        errors.append("stats warnings missing or malformed")
    else:
        flagged = [
            warning
            for warning in warnings
            if warning in STRICT_STATS and not waived(evidence, "stats", warning=warning)
        ]
        if flagged:
            errors.append("stats warnings require remediation or an explicit waiver: " + ", ".join(flagged))


def check_content(
    evidence: dict[str, Any], root: Path, expected_slides: set[int], errors: list[str]
) -> None:
    interview = evidence.get("interview")
    if not isinstance(interview, dict):
        errors.append("interview evidence missing")
    else:
        if interview.get("mode") not in {"answered", "auto"}:
            errors.append("interview.mode must be answered or auto")
        require_string(interview.get("record"), "interview.record", errors, minimum=12)

    content = evidence.get("content")
    if not isinstance(content, dict):
        errors.append("content evidence missing")
        return
    source_mode = content.get("source_mode")
    if source_mode not in {"provided", "web", "none"}:
        errors.append("content.source_mode must be provided, web, or none")
        source_mode = "provided"
    sources = content.get("sources")
    if not isinstance(sources, list) or (source_mode != "none" and not sources):
        errors.append("content.sources must document the supplied source material")
    elif isinstance(sources, list):
        for index, source in enumerate(sources, start=1):
            label = f"content.sources[{index}]"
            if not isinstance(source, dict):
                errors.append(f"{label} must be an object")
                continue
            kind = source.get("kind")
            if kind == "provided":
                check_hashed_file(root, source.get("path"), source.get("sha256"), label, errors)
            elif kind == "web":
                require_string(source.get("locator"), f"{label}.locator", errors, minimum=12)
            else:
                errors.append(f"{label}.kind must be provided or web")

    slides = content.get("slides")
    if not isinstance(slides, list):
        errors.append("content.slides missing or malformed")
    else:
        rows = {row.get("slide"): row for row in slides if isinstance(row, dict) and isinstance(row.get("slide"), int)}
        if set(rows) != expected_slides:
            errors.append("content.slides must cover every final slide exactly once")
        for number, row in rows.items():
            require_string(row.get("role"), f"content slide {number}.role", errors)
            require_string(row.get("takeaway"), f"content slide {number}.takeaway", errors, minimum=8)
            evidence_rows = row.get("evidence")
            if not isinstance(evidence_rows, list) or not all(isinstance(item, str) and item.strip() for item in evidence_rows):
                errors.append(f"content slide {number}.evidence must contain source references")

    ledger = content.get("claim_ledger")
    if not isinstance(ledger, list) or (source_mode != "none" and not ledger):
        errors.append("content.claim_ledger must contain verified claims")
    elif isinstance(ledger, list):
        for index, claim in enumerate(ledger, start=1):
            label = f"content.claim_ledger[{index}]"
            if not isinstance(claim, dict):
                errors.append(f"{label} must be an object")
                continue
            require_string(claim.get("claim"), f"{label}.claim", errors, minimum=8)
            require_string(claim.get("source"), f"{label}.source", errors, minimum=3)
            if claim.get("verified") is not True:
                errors.append(f"{label}.verified must be true")

    checkpoint = content.get("checkpoint")
    if not isinstance(checkpoint, dict):
        errors.append("content.checkpoint missing")
    else:
        if checkpoint.get("mode") not in {"approved", "auto"}:
            errors.append("content.checkpoint.mode must be approved or auto")
        require_string(checkpoint.get("record"), "content.checkpoint.record", errors, minimum=12)


def check_design(
    evidence: dict[str, Any],
    root: Path,
    expected_slides: set[int],
    deck_hash: str,
    errors: list[str],
) -> dict[int, dict[str, Any]]:
    design = evidence.get("design")
    if not isinstance(design, dict):
        errors.append("design evidence missing")
        return {}
    direction = design.get("direction")
    if not isinstance(direction, dict):
        errors.append("design.direction missing")
    else:
        branch = direction.get("branch")
        if branch not in {"clean", "provided-template", "generated-template", "mimic"}:
            errors.append("design.direction.branch is invalid")
        artifact_path = check_hashed_file(
            root,
            direction.get("artifact"),
            direction.get("sha256"),
            "design.direction.artifact",
            errors,
            minimum_bytes=512,
        )
        if artifact_path is not None and artifact_path.suffix.lower() in {".html", ".htm"}:
            markup = artifact_path.read_text(encoding="utf-8", errors="ignore").lower()
            if "<html" not in markup or "<body" not in markup:
                errors.append("design.direction.artifact must be a real HTML preview, not a placeholder file")
        choices = direction.get("directions")
        if not isinstance(choices, list) or not choices:
            errors.append("design.direction.directions missing")
        elif branch == "clean":
            names = [choice.get("name") for choice in choices if isinstance(choice, dict)]
            if len(choices) < 4 or len(set(names)) < 4 or any(not isinstance(name, str) or not name.strip() for name in names):
                errors.append("clean design direction needs four named preview directions")
            else:
                try:
                    validator = load_direction_validator()
                    result = validator.check(choices)
                    if result.get("flagged") and not require_string(
                        direction.get("diversity_waiver"), "design.direction.diversity_waiver", errors, minimum=12
                    ):
                        errors.append("direction preview has too-similar candidates without a named diversity waiver")
                except Exception as exc:
                    errors.append(f"design.direction candidates cannot pass the diversity check: {exc}")
            if artifact_path is not None and artifact_path.suffix.lower() in {".html", ".htm"}:
                markup = artifact_path.read_text(encoding="utf-8", errors="ignore")
                missing_names = [name for name in names if isinstance(name, str) and name not in markup]
                if missing_names:
                    errors.append("direction preview does not visibly include: " + ", ".join(missing_names))
        if direction.get("decision") not in {"user-approved", "auto-carve", "provided-template", "mode-a-mimic"}:
            errors.append("design.direction.decision is invalid")
        require_string(direction.get("record"), "design.direction.record", errors, minimum=12)

    scale = design.get("type_scale")
    if not isinstance(scale, dict):
        errors.append("design.type_scale missing")
    else:
        for key in ("display", "title", "body"):
            if not isinstance(scale.get(key), (int, float)):
                errors.append(f"design.type_scale.{key} must be numeric")
        if isinstance(scale.get("body"), (int, float)) and scale["body"] < BODY_FLOORS[evidence.get("delivery", "presented")]:
            errors.append("design.type_scale.body is below the delivery floor")

    # The documented dial is <conservative | balanced+ | bold | experimental> (SKILL.md,
    # agents/slide-design.md, review-rubrics.md). This set had `conservative` MISSING and
    # `deliberately-restrained` in its place — but that string is a value for the signature_move
    # FIELD ("deliberately restrained: <why>", the conservative dial's documented escape), never a
    # dial. So the one word an author is actually told to write failed this gate, and a field value
    # passed as a dial. The legacy string stays accepted so existing evidence files still load.
    if design.get("boldness") not in {"conservative", "balanced+", "bold", "experimental",
                                      "deliberately-restrained"}:
        errors.append("design.boldness must declare a supported direction "
                      "(conservative | balanced+ | bold | experimental)")
    require_string(design.get("signature_move"), "design.signature_move", errors, minimum=12)
    carried_by = design.get("carried_by")
    if not isinstance(carried_by, list) or len(set(carried_by) & expected_slides) < min(2, len(expected_slides)):
        errors.append("design.carried_by must show the signature move on at least two slides")

    # Mirrors render_deck.py's carve with the SAME condition — the two gate paths disagreeing about
    # what an honest plan looks like has already cost this repo once (this key was spelled `path`
    # here and `png` there, so a bridged run wrote the field its own gate demanded and the other
    # rejected it).
    carved = (str(design.get("boldness", "")).strip().lower() == "conservative"
              and str(design.get("signature_move", "")).strip().lower()
                  .startswith("deliberately restrained"))
    proof = design.get("signature_proof")
    if carved and proof is None:
        pass                          # a conservative deck that declared it took no risk
    elif not isinstance(proof, dict):
        errors.append("design.signature_proof missing")
    else:
        proof_slide = proof.get("slide")
        if proof_slide not in expected_slides:
            errors.append("design.signature_proof.slide is not a final slide")
        proof_path = check_hashed_file(
            root,
            proof.get("path"),
            proof.get("sha256"),
            "design.signature_proof",
            errors,
            minimum_bytes=512,
        )
        if proof_path is not None:
            dimensions = png_dimensions(proof_path)
            if dimensions is None or dimensions[0] < 640 or dimensions[1] < 360:
                errors.append("design.signature_proof must be a rendered PNG of at least 640x360")
        if proof.get("pptx_sha256") != deck_hash:
            errors.append("signature proof is not bound to the final PPTX SHA-256")

    rows = design.get("slides")
    if not isinstance(rows, list):
        errors.append("design.slides missing or malformed")
        row_map: dict[int, dict[str, Any]] = {}
    else:
        row_map = {row.get("slide"): row for row in rows if isinstance(row, dict) and isinstance(row.get("slide"), int)}
        if set(row_map) != expected_slides:
            errors.append("design.slides must cover every final slide exactly once")
        for number, row in row_map.items():
            require_string(row.get("function"), f"design slide {number}.function", errors)
            form = require_string(row.get("form"), f"design slide {number}.form", errors)
            runner_up = require_string(row.get("runner_up"), f"design slide {number}.runner_up", errors)
            if form and runner_up and form.lower() == runner_up.lower():
                errors.append(f"design slide {number} must keep a distinct runner-up form")
            require_string(row.get("reason"), f"design slide {number}.reason", errors, minimum=12)
            if not isinstance(row.get("categorical"), bool):
                errors.append(f"design slide {number}.categorical must be boolean")
            components = row.get("components")
            if not isinstance(components, list) or not all(isinstance(component, str) and component for component in components):
                errors.append(f"design slide {number}.components must be a list of component names")

    checkpoint = design.get("checkpoint")
    if not isinstance(checkpoint, dict):
        errors.append("design.checkpoint missing")
    else:
        if checkpoint.get("mode") not in {"approved", "auto"}:
            errors.append("design.checkpoint.mode must be approved or auto")
        require_string(checkpoint.get("record"), "design.checkpoint.record", errors, minimum=12)
    return row_map


def check_build(
    evidence: dict[str, Any], root: Path, supplied_script: Path, errors: list[str]
) -> tuple[Path | None, dict[str, set[str]]]:
    build = evidence.get("build")
    if not isinstance(build, dict):
        errors.append("build evidence missing")
        return None, {}
    script = check_hashed_file(root, build.get("script"), build.get("sha256"), "build.script", errors)
    if script is None:
        return None, {}
    if script.resolve() != supplied_script.resolve():
        errors.append("--build-script does not match evidence.build.script")
    if build.get("strict_layout") is not True:
        errors.append("build.strict_layout must be true")
    calls, has_strict_layout = parse_script(script, errors)
    for token in forbidden_icon_rasterizer_calls(script, errors):
        errors.append(
            f"build script invokes macOS Quick Look thumbnail generation ({token}); use scripts/icons.py / icon_png for transparent high-resolution icon assets"
        )
    if not has_strict_layout:
        errors.append("build script must call lint_layout(..., strict=True)")
    return script, calls


def check_components(
    evidence: dict[str, Any],
    components: dict[str, Any],
    design_rows: dict[int, dict[str, Any]],
    calls: dict[str, set[str]],
    errors: list[str],
) -> None:
    for slide, row in design_rows.items():
        function = row.get("function")
        if not isinstance(function, str) or function not in calls:
            errors.append(f"design slide {slide} function is absent from the build script")
            continue
        missing = [component for component in row.get("components", []) if component not in calls[function]]
        if missing:
            errors.append(f"design slide {slide} declares component(s) not called by {function}: {', '.join(missing)}")

    clusters = components.get("clusters", [])
    if not isinstance(clusters, list):
        errors.append("component audit clusters missing or malformed")
        return
    for cluster in clusters:
        if not isinstance(cluster, dict):
            continue
        slide = cluster.get("slide")
        suggestions = cluster.get("suggest", [])
        if slide not in design_rows or not isinstance(suggestions, list) or not suggestions:
            continue
        row = design_rows[slide]
        function = row.get("function")
        implemented = set(row.get("components", [])) & calls.get(function, set()) & set(suggestions)
        suppressed_by = {
            item for item in components.get("suppressed_by", []) if isinstance(item, str)
        }
        emitted_here = suppressed_by & calls.get(function, set()) & set(row.get("components", []))
        if not implemented and not emitted_here and not waived(
            evidence, "component", slide=slide, pattern=cluster.get("pattern")
        ):
            pattern = cluster.get("pattern", "unnamed cluster")
            errors.append(
                f"component audit cluster on slide {slide} ({pattern}) needs a suggested component or a documented waiver"
            )


def recompute_component_audit(script: Path, deck: Path, errors: list[str]) -> dict[str, Any] | None:
    command = [
        sys.executable,
        str(Path(__file__).with_name("component_audit.py")),
        str(script),
        str(deck),
        "--json",
    ]
    result = subprocess.run(command, capture_output=True, text=True)
    if result.returncode not in {0, 2}:
        detail = (result.stdout + result.stderr).strip().replace("\n", " ")
        errors.append("component audit could not inspect the final build/deck: " + detail[:300])
        return None
    try:
        value = json.loads(result.stdout)
    except json.JSONDecodeError as exc:
        errors.append(f"component audit returned invalid JSON: {exc}")
        return None
    if not isinstance(value, dict) or value.get("inspected") is not True:
        errors.append("component audit did not confirm inspection of the final deck")
        return None
    return value


def check_icons(
    evidence: dict[str, Any],
    root: Path,
    design_rows: dict[int, dict[str, Any]],
    calls: dict[str, set[str]],
    errors: list[str],
) -> None:
    icon_rows = evidence.get("icons", [])
    if not isinstance(icon_rows, list):
        errors.append("icons must be a list")
        return
    by_slide: dict[int, list[dict[str, Any]]] = defaultdict(list)
    families: set[str] = set()
    for index, row in enumerate(icon_rows, start=1):
        label = f"icons[{index}]"
        if not isinstance(row, dict) or not isinstance(row.get("slide"), int):
            errors.append(f"{label} must identify a slide")
            continue
        family = require_string(row.get("family"), f"{label}.family", errors)
        if family:
            families.add(family)
        asset = check_hashed_file(root, row.get("asset"), row.get("sha256"), label, errors, minimum_bytes=32)
        rasterizer = require_string(row.get("rasterizer"), f"{label}.rasterizer", errors, minimum=3)
        if rasterizer not in {"scripts/icons.py", "provided-hires"}:
            errors.append(f"{label}.rasterizer must be scripts/icons.py or provided-hires")
        if asset is not None:
            info = png_info(asset)
            if info is None:
                errors.append(f"{label} must be a readable PNG icon asset")
            else:
                width, height, _bit_depth, color_type = info
                if min(width, height) < 256:
                    errors.append(
                        f"{label} is only {width}x{height}px; Codex icon assets need a 256px minimum edge to avoid thumbnail blur"
                    )
                if color_type not in {4, 6}:
                    errors.append(
                        f"{label} must preserve transparent alpha (PNG colour type 4 or 6), not a matted thumbnail"
                    )
        by_slide[row["slide"]].append(row)
    if len(families) > 1:
        errors.append("all icon evidence must use one coherent icon family")

    for slide, design in design_rows.items():
        if not design.get("categorical"):
            continue
        rows = by_slide.get(slide, [])
        if not rows:
            if not waived(evidence, "icon", slide=slide):
                errors.append(f"categorical slide {slide} needs rendered icon evidence or a documented waiver")
            continue
        function = design.get("function")
        helper_count = len(calls.get(function, set()) & ICON_HELPERS)
        if helper_count < 1:
            errors.append(f"categorical slide {slide} records icon assets but {function} has no deckkit icon helper call")


def load_review_validator() -> Any:
    validator_path = Path(__file__).with_name("validate_review.py")
    spec = importlib.util.spec_from_file_location("slide_maker_validate_review", validator_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load validate_review.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_direction_validator() -> Any:
    validator_path = Path(__file__).with_name("directions_diversity.py")
    spec = importlib.util.spec_from_file_location("slide_maker_directions_diversity", validator_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load directions_diversity.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_visual_contract() -> Any:
    script_path = Path(__file__).with_name("codex_visual_contract.py")
    spec = importlib.util.spec_from_file_location("slide_maker_codex_visual_contract", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("could not load Codex visual-contract checker")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def check_critics(
    evidence: dict[str, Any], root: Path, expected_slides: set[int], deck_hash: str, errors: list[str]
) -> dict[str, dict[str, Any]]:
    reviewed: dict[str, dict[str, Any]] = {}
    critics = evidence.get("critics")
    if not isinstance(critics, list):
        errors.append("critics must be a list")
        return reviewed
    effort = evidence.get("review_effort", "standard")
    if effort not in {"fast", "standard", "thorough"}:
        errors.append("review_effort must be fast, standard, or thorough")
        return reviewed
    lenses = [row.get("lens") for row in critics if isinstance(row, dict)]
    thorough_scope = None
    if effort == "fast":
        if not require_string(evidence.get("fast_opt_in"), "fast_opt_in", errors, minimum=12):
            return reviewed
        required_lenses = {"general"}
    else:
        required_lenses = {"content", "design"}
    if effort == "thorough":
        panel = evidence.get("thorough_panel")
        if not isinstance(panel, dict):
            errors.append("thorough review requires thorough_panel evidence")
        else:
            thorough_scope = panel.get("scope")
            if thorough_scope not in {"light", "full"}:
                errors.append("thorough_panel.scope must be light or full")
            require_string(panel.get("record"), "thorough_panel.record", errors, minimum=12)
    if not required_lenses.issubset(set(lenses)):
        errors.append("critic panel does not cover the required lenses")
    if len(set(row.get("review") for row in critics if isinstance(row, dict))) != len(critics):
        errors.append("each critic lens must use a distinct review artifact")

    try:
        validator = load_review_validator()
    except Exception as exc:
        errors.append(f"could not load critic validator: {exc}")
        return
    for index, critic in enumerate(critics, start=1):
        label = f"critics[{index}]"
        if not isinstance(critic, dict):
            errors.append(f"{label} must be an object")
            continue
        lens = critic.get("lens")
        if lens not in {"content", "design", "general"}:
            errors.append(f"{label}.lens is invalid")
            continue
        path = check_hashed_file(root, critic.get("review"), critic.get("sha256"), label, errors, minimum_bytes=32)
        if path is None:
            continue
        if critic.get("pptx_sha256") != deck_hash:
            errors.append(f"{label} is not bound to the final PPTX SHA-256")
        try:
            review = load_json(path)
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            errors.append(f"{label} cannot be read: {exc}")
            continue
        reviewer = review.get("reviewer")
        if not isinstance(reviewer, dict):
            errors.append(f"{label} must record reviewer provenance (isolated, human, or self-review)")
        else:
            origin = reviewer.get("origin")
            if origin not in {"isolated", "human", "self-review"}:
                errors.append(f"{label}.reviewer.origin must be isolated, human, or self-review")
            require_string(reviewer.get("identity"), f"{label}.reviewer.identity", errors, minimum=3)
            if reviewer.get("fresh_context") is not True:
                errors.append(f"{label}.reviewer.fresh_context must be true")
            if origin == "self-review" and not waived(evidence, "critic-independence", lens=lens):
                errors.append(f"{label} is self-review; record a critic-independence waiver instead of claiming independent consent")
        for issue in validator.validate_critic(review):
            errors.append(f"{label} invalid critic schema: {issue}")
        coverage = review.get("coverage", {})
        opened = coverage.get("slides_opened", []) if isinstance(coverage, dict) else []
        if not isinstance(opened, list) or not expected_slides.issubset(set(opened)):
            errors.append(f"{label} did not inspect every final slide")
        if not isinstance(coverage, dict) or coverage.get("stats_block_seen") is not True:
            errors.append(f"{label} did not confirm the statistics block")
        if not isinstance(coverage, dict) or coverage.get("contract_card_seen") is not True:
            errors.append(f"{label} did not confirm the design contract")
        passes = coverage.get("passes", []) if isinstance(coverage, dict) else []
        if not isinstance(passes, list) or not any(lens in str(item).lower() for item in passes):
            errors.append(f"{label} does not record its {lens} lens pass")
        if lens == "general":
            normalized = " ".join(str(item).lower() for item in passes)
            if "content" not in normalized or "design" not in normalized:
                errors.append(f"{label} generalist pass must explicitly cover content and design")
        probes = review.get("probes", {})
        if lens == "content" and not isinstance(probes, dict):
            errors.append(f"{label} lacks content probes")
        if lens == "content" and not require_string(probes.get("memory_sentence") if isinstance(probes, dict) else None, f"{label}.probes.memory_sentence", errors, minimum=12):
            pass
        if lens == "design":
            per_slide = probes.get("per_slide", []) if isinstance(probes, dict) else []
            inspected = {row.get("slide") for row in per_slide if isinstance(row, dict)}
            if not expected_slides.issubset(inspected):
                errors.append(f"{label} lacks a design probe for every final slide")
        reviewed[lens] = review

    if effort == "thorough" and thorough_scope == "full":
        arbiters = evidence.get("arbiters")
        if not isinstance(arbiters, list) or not arbiters:
            errors.append("full thorough review requires a final arbiter confirmation artifact")
            return reviewed
        for index, arbiter in enumerate(arbiters, start=1):
            label = f"arbiters[{index}]"
            if not isinstance(arbiter, dict):
                errors.append(f"{label} must be an object")
                continue
            path = check_hashed_file(root, arbiter.get("review"), arbiter.get("sha256"), label, errors, minimum_bytes=32)
            if path is None:
                continue
            if arbiter.get("pptx_sha256") != deck_hash:
                errors.append(f"{label} is not bound to the final PPTX SHA-256")
            try:
                review = load_json(path)
            except (OSError, ValueError, json.JSONDecodeError) as exc:
                errors.append(f"{label} cannot be read: {exc}")
                continue
            for issue in validator.validate_arbiter(review):
                errors.append(f"{label} invalid arbiter schema: {issue}")
            checks = review.get("checks")
            if not isinstance(checks, list):
                errors.append(f"{label} must contain final fix-confirmation checks")
            elif any(not row.get("resolved") or row.get("dulled") for row in checks if isinstance(row, dict)):
                errors.append(f"{label} reports an unresolved or dulled final fix")
    return reviewed


def check_visual_contract(
    evidence: dict[str, Any],
    root: Path,
    deck_path: Path | None,
    deck_hash: str,
    build_script: Path | None,
    reviews: dict[str, dict[str, Any]],
    errors: list[str],
) -> None:
    contract = evidence.get("visual_contract")
    if not isinstance(contract, dict):
        errors.append("visual_contract evidence is required for Codex delivery")
        return
    manifest_path = check_hashed_file(
        root, contract.get("manifest"), contract.get("sha256"), "visual_contract.manifest", errors, minimum_bytes=32
    )
    result_path = check_hashed_file(
        root, contract.get("result"), contract.get("result_sha256"), "visual_contract.result", errors, minimum_bytes=32
    )
    if contract.get("pptx_sha256") != deck_hash:
        errors.append("visual_contract is not bound to the final PPTX SHA-256")
    if manifest_path is None or result_path is None or deck_path is None or build_script is None:
        return
    try:
        manifest = load_json(manifest_path)
        result = load_json(result_path)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        errors.append(f"visual_contract cannot be read: {exc}")
        return
    if result.get("schema") != "slide-maker-codex-visual-contract-result/v1":
        errors.append("visual_contract result has the wrong schema")
    if result.get("pptx_sha256") != deck_hash:
        errors.append("visual_contract result was produced for a different PPTX")
    if result.get("manifest_sha256") != contract.get("sha256"):
        errors.append("visual_contract result is not bound to the supplied manifest")
    if result.get("passed") is not True:
        errors.append("visual_contract result contains failed local checks")
    try:
        checker = load_visual_contract()
        fresh = checker.evaluate(deck_path, manifest_path, build_script)
    except Exception as exc:
        errors.append(f"visual_contract could not be recomputed: {exc}")
        return
    if fresh.get("passed") is not True:
        errors.append("visual_contract fails when recomputed against the final build and PPTX")
    declared_zones = {
        row.get("id") for row in manifest.get("zones", []) if isinstance(row, dict) and isinstance(row.get("id"), str)
    }
    declared_icons = {
        row.get("id") for row in manifest.get("icons", []) if isinstance(row, dict) and isinstance(row.get("id"), str)
    }
    result_zones = {
        row.get("id") for row in result.get("zones", []) if isinstance(row, dict) and row.get("pass") is True
    }
    result_icons = {
        row.get("id") for row in result.get("icons", []) if isinstance(row, dict) and row.get("pass") is True
    }
    if result_zones != declared_zones or result_icons != declared_icons:
        errors.append("visual_contract result does not account for every declared zone and icon")
    design_review = reviews.get("design")
    probes = design_review.get("probes", {}) if isinstance(design_review, dict) else {}
    hotspot_checks = probes.get("hotspot_checks", []) if isinstance(probes, dict) else []
    icon_checks = probes.get("icon_checks", []) if isinstance(probes, dict) else []
    checked_zones = {
        row.get("id") for row in hotspot_checks
        if isinstance(row, dict) and row.get("result") == "pass" and isinstance(row.get("observed"), str) and len(row["observed"].strip()) >= 12
    }
    checked_icons = {
        row.get("id") for row in icon_checks
        if isinstance(row, dict) and row.get("result") == "pass" and isinstance(row.get("observed"), str) and len(row["observed"].strip()) >= 12
    }
    if not declared_zones.issubset(checked_zones):
        errors.append("design critic did not visually attest every declared local hotspot")
    if not declared_icons.issubset(checked_icons):
        errors.append("design critic did not attest every declared icon's semantic fit")


def evaluate(
    lint: dict[str, Any], components: dict[str, Any], supplied_script: Path, evidence: dict[str, Any], root: Path
) -> list[str]:
    errors: list[str] = []
    if evidence.get("schema") != SCHEMA:
        errors.append(f"evidence.schema must be {SCHEMA}")
    if evidence.get("runtime") not in {"codex", "openai-gpt-bridged"}:
        errors.append("evidence.runtime must be codex or openai-gpt-bridged")
    delivery = evidence.get("delivery", "presented")
    if delivery not in BODY_FLOORS:
        errors.append("evidence.delivery must be presented, textheavy, or selfread")
        delivery = "presented"
    deck = evidence.get("deck")
    deck_hash = ""
    expected_slides: set[int] = set()
    deck_path: Path | None = None
    if not isinstance(deck, dict):
        errors.append("deck evidence missing")
    else:
        deck_path = check_hashed_file(root, deck.get("pptx"), deck.get("sha256"), "deck", errors, minimum_bytes=512)
        deck_hash = deck.get("sha256") if is_sha256(deck.get("sha256")) else ""
        count = deck.get("slide_count")
        if not isinstance(count, int) or count < 1:
            errors.append("deck.slide_count must be a positive integer")
        else:
            expected_slides = set(range(1, count + 1))
            if deck_path is not None:
                actual_count = slide_count_from_pptx(deck_path)
                if actual_count is None:
                    errors.append("could not read final PPTX to count slides")
                elif actual_count != count:
                    errors.append(f"deck.slide_count is {count}, but final PPTX contains {actual_count} slides")
    check_lint(lint, delivery, evidence, errors)
    if expected_slides:
        check_content(evidence, root, expected_slides, errors)
        design_rows = check_design(evidence, root, expected_slides, deck_hash, errors)
        build_script, calls = check_build(evidence, root, supplied_script, errors)
        audited_components = components
        if build_script is not None and deck_path is not None:
            recomputed = recompute_component_audit(build_script, deck_path, errors)
            if recomputed is not None:
                audit_keys = ("used_forms", "clusters", "suppressed_by", "inspected")
                if any(components.get(key) != recomputed.get(key) for key in audit_keys):
                    errors.append("components JSON does not match a fresh audit of the final build and PPTX")
                audited_components = recomputed
        check_components(evidence, audited_components, design_rows, calls, errors)
        check_icons(evidence, root, design_rows, calls, errors)
        reviews = check_critics(evidence, root, expected_slides, deck_hash, errors)
        check_visual_contract(evidence, root, deck_path, deck_hash, build_script, reviews, errors)
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description="Codex-only final delivery gate for slide-maker")
    parser.add_argument("--lint", type=Path, help="JSON from lint_layout")
    parser.add_argument("--components", type=Path, help="JSON from component_audit.py")
    parser.add_argument("--build-script", type=Path, help="final deck build script")
    parser.add_argument("--evidence", type=Path, help=".codex-deck-evidence.json")
    parser.add_argument("--init", type=Path, help="write an evidence template and exit")
    args = parser.parse_args()

    if args.init:
        if args.init.exists():
            print(f"refusing to overwrite existing file: {args.init}", file=sys.stderr)
            return 2
        args.init.write_text(json.dumps(TEMPLATE, indent=2) + "\n", encoding="utf-8")
        print(f"wrote evidence template: {args.init}")
        return 0

    missing = [name for name in ("lint", "components", "build_script", "evidence") if getattr(args, name) is None]
    if missing:
        parser.error("required unless --init: " + ", ".join("--" + item.replace("_", "-") for item in missing))
    try:
        lint = load_json(args.lint)
        components = load_json(args.components)
        evidence = load_json(args.evidence)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"cannot read gate input: {exc}", file=sys.stderr)
        return 2
    errors = evaluate(lint, components, args.build_script, evidence, args.evidence.parent)
    if errors:
        print("CODEX DELIVERY GATE: BLOCKED")
        for error in errors:
            print(f"- {error}")
        return 1
    print("CODEX DELIVERY GATE: PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
