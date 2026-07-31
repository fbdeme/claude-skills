#!/usr/bin/env python3
"""Render a .pptx to one PNG per slide, so you can SEE each slide and catch
overflow / contrast / glyph problems before handing the deck back.

Cross-platform: works on macOS, Linux, WSL, and NATIVE Windows (PowerShell / cmd) —
no shell required. The .sh wrapper just delegates here.

Usage:
    python render_deck.py /path/to/deck.pptx [out_dir]
    # Windows:  python scripts\\render_deck.py C:\\path\\deck.pptx
Output: <out_dir>/slide01.png, slide02.png, ...   (default out_dir: ./render)

Requires: LibreOffice + pymupdf (python -m pip install pymupdf). One-time installs.
Override LibreOffice discovery with the SOFFICE env var (full path to the binary).
"""
import json
import os
import re
import sys
import shutil
import tempfile
import subprocess
from pathlib import Path


def find_soffice():
    """Locate the LibreOffice binary across macOS / Linux / WSL / native Windows.

    Order: $SOFFICE override -> anything on PATH -> known install locations.
    On Windows, prefer the sibling ``soffice.com`` (the console front-end that
    blocks until conversion finishes) over ``soffice.exe`` (which can detach and
    leave the PDF half-written).
    """
    def prefer_com(path):
        if path and path.lower().endswith("soffice.exe"):
            com = path[:-4] + ".com"
            if os.path.isfile(com):
                return com
        return path

    env = os.environ.get("SOFFICE")
    if env and os.path.isfile(env):
        return prefer_com(env)

    for cmd in ("soffice", "libreoffice", "soffice.com", "soffice.exe"):
        found = shutil.which(cmd)
        if found:
            return prefer_com(found)

    candidates = [
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        # Linux
        "/usr/bin/soffice", "/usr/bin/libreoffice", "/usr/local/bin/soffice",
        "/snap/bin/libreoffice", "/opt/libreoffice/program/soffice",
        # native Windows (default installer locations)
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # Windows install reached from WSL via the /mnt/c mount
        "/mnt/c/Program Files/LibreOffice/program/soffice.exe",
    ]
    # %ProgramFiles% may point somewhere non-default on Windows.
    for var in ("ProgramFiles", "ProgramFiles(x86)", "ProgramW6432"):
        base = os.environ.get(var)
        if base:
            candidates.append(os.path.join(base, "LibreOffice", "program", "soffice.exe"))

    for path in candidates:
        if os.path.isfile(path):
            return prefer_com(path)
    return None


def die(msg, code=1):
    print(msg, file=sys.stderr)
    sys.exit(code)


def _tail(text, limit=4000):
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return "...<truncated>...\n" + text[-limit:]


# Exactly what a render writes: slide01.png … slideNN.png and the two bookend thumbnails. Used to
# decide what may be deleted when the out dir is shared with the user's own files.
_RENDER_PNG = re.compile(r"^(slide\d{2,}|thumb_first|thumb_last)\.png$")


def _rels_targets(xml_bytes, base_dir):
    """Parse a .rels part -> {rId: normalized package path}.

    Parsed as XML, not by regex: OOXML allows the Id/Target attributes in either order, and a
    Target may be absolute ("/ppt/slides/slide1.xml"). Getting this wrong is not a cosmetic bug —
    a rId that fails to resolve silently drops a slide from the deck order, which shifts every
    index after it and writes real slides into the wrong PNG filenames.
    """
    import posixpath
    import xml.etree.ElementTree as ET
    out = {}
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return out
    for rel in root:
        rid = rel.get("Id")
        tgt = rel.get("Target")
        if not rid or not tgt or (rel.get("TargetMode") or "") == "External":
            continue
        if tgt.startswith("/"):
            out[rid] = tgt.lstrip("/")
        else:
            out[rid] = posixpath.normpath(posixpath.join(base_dir, tgt))
    return out


def _slide_fingerprints(pptx):
    """One content hash per slide, in deck order, plus any reason the deck cannot be diffed.

    Each slide's hash covers its own XML, its rels, the BYTES of every media part it references,
    and a DECK-GLOBAL digest (presentation.xml, theme, masters, layouts and the media those
    reference) — so a re-cropped photo, a swapped plate, a new theme or a different canvas size all
    count as changes.

    Returns (hashes, blockers). `blockers` is a list of human-readable reasons the caller must fall
    back to a full render; it is never a soft signal. Anything that makes the slide->page mapping
    uncertain belongs here, because a wrong mapping writes a real slide into another slide's PNG.
    """
    import hashlib
    import zipfile
    blockers = []
    with zipfile.ZipFile(pptx) as z:
        names = set(z.namelist())
        _part_cache = {}

        def part_digest(name):
            """Hash a package part ONCE. A background plate shared by every slide was previously
            re-read and re-hashed per referencing slide, making cost O(slides x media-bytes)."""
            d = _part_cache.get(name)
            if d is None:
                d = hashlib.sha256(z.read(name)).digest()
                _part_cache[name] = d
            return d
        pres_rels = _rels_targets(z.read("ppt/_rels/presentation.xml.rels"), "ppt") \
            if "ppt/_rels/presentation.xml.rels" in names else {}

        # ---- deck-global digest -------------------------------------------------------
        # docProps/* is excluded on purpose: core.xml carries a modified-timestamp that changes on
        # every save, which would force a full render every time and delete the feature.
        # ppt/media and ppt/notesSlides are excluded here — slide media is covered per-slide below,
        # and notes never reach a rendered pixel. But media referenced by a LAYOUT, MASTER or THEME
        # is deck-global and IS folded in: swapping a master's background image in place changes no
        # slide's XML, so without this the whole deck would render stale under "no slide changed".
        gh = hashlib.sha256()
        global_media = set()
        for n in sorted(names):
            if not n.startswith("ppt/"):
                continue
            if n.startswith(("ppt/slides/", "ppt/media/", "ppt/notesSlides/")):
                continue
            gh.update(n.encode())
            gh.update(z.read(n))
            if n.startswith(("ppt/slideLayouts/_rels/", "ppt/slideMasters/_rels/", "ppt/theme/_rels/")):
                base = n.rsplit("/_rels/", 1)[0]
                for tgt in _rels_targets(z.read(n), base).values():
                    if "/media/" in "/" + tgt:
                        global_media.add(tgt)
        for m in sorted(global_media):
            if m in names:
                gh.update(m.encode())
                gh.update(part_digest(m))
        global_digest = gh.digest()

        # ---- deck order ---------------------------------------------------------------
        import re
        pres = z.read("ppt/presentation.xml").decode("utf-8", "replace")
        rids = re.findall(r'<p:sldId[^>]*r:id="([^"]+)"', pres)
        order = [pres_rels.get(r) for r in rids]
        if any(t is None for t in order):
            blockers.append("a slide reference could not be resolved in presentation.xml.rels")
            order = [t for t in order if t]

        fps = []
        for sname in order:
            if sname not in names:
                # Never hash a missing part as b"" — every affected slide would collapse to the
                # same constant and real edits would stop registering, forever.
                blockers.append("slide part missing from the package: {}".format(sname))
                continue
            h = hashlib.sha256()
            h.update(global_digest)
            body = z.read(sname)
            h.update(body)
            if b'type="slidenum"' in body:
                # An auto slide-number field renumbers itself inside a subset.
                blockers.append("deck uses auto slide-number fields")
            if re.search(rb'<p:sld\b[^>]*\bshow="0"', body):
                # LibreOffice DROPS hidden slides from the PDF, so PDF page N is no longer deck
                # slide N — on the full path too. Verified with a 4-slide deck: one hidden slide
                # produced a 3-page PDF.
                blockers.append("deck has hidden slides")
            rname = sname.replace("slides/", "slides/_rels/") + ".rels"
            rel_bytes = z.read(rname) if rname in names else b""
            h.update(rel_bytes)
            for tgt in sorted(_rels_targets(rel_bytes, "ppt/slides").values()):
                if "/media/" in "/" + tgt and tgt in names:
                    h.update(tgt.encode())
                    h.update(part_digest(tgt))
            fps.append(h.hexdigest())

        if len(fps) != len(rids):
            blockers.append("deck order could not be fully resolved ({} of {} slides)".format(
                len(fps), len(rids)))
    # de-dup, order-stable
    seen, uniq = set(), []
    for b in blockers:
        if b not in seen:
            seen.add(b); uniq.append(b)
    return fps, uniq


def _subset_pptx(src, keep_idx, dest):
    """Copy `src` keeping ONLY the 0-indexed slides in `keep_idx` (order preserved).

    Deleting sldId entries is safe and cheap; the orphaned parts stay in the package and
    LibreOffice ignores them. Verified pixel-identical to the same slide rendered from the
    full deck, which is what makes the fast path trustworthy rather than merely fast.
    """
    import shutil
    from pptx import Presentation
    shutil.copy(src, dest)
    prs = Presentation(dest)
    lst = prs.slides._sldIdLst
    for i, sid in enumerate(list(lst)):
        if i not in keep_idx:
            lst.remove(sid)
    prs.save(dest)
    return dest


def _render_pdf(soffice, src, outdir):
    """pptx -> pdf via a throwaway LibreOffice profile, into an EMPTY private directory.

    `outdir` must not already contain `<src-stem>.pdf`. Rendering into a directory that may hold a
    previous PDF is how a FAILED conversion gets read back as success: the caller checks only that
    the file exists, and a stale one satisfies that. Reproduced with real LibreOffice on a deck it
    refuses to convert — the run printed "rendered N slides" and exit 0 over untouched output.
    Returns (pdf_path, result, cmd); the caller must check `result.returncode`.
    """
    profile = tempfile.mkdtemp(prefix="lo_render_")
    try:
        cmd = [soffice, "-env:UserInstallation=" + Path(profile).as_uri(),
               "--headless", "--convert-to", "pdf", "--outdir", outdir, src]
        try:
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                                    text=True, timeout=300)
        except subprocess.TimeoutExpired:
            die("LibreOffice render exceeded 300s and was killed — the .pptx may be malformed or "
                "hostile (e.g. a decompression bomb). Nothing was produced from {}.".format(src))
    finally:
        shutil.rmtree(profile, ignore_errors=True)
    return (os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + ".pdf"), result, cmd)


def _sha256(path):
    import hashlib
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


GATES_FILE = ".deck-gates.json"


def check_handoff_gates(pptx, mode="presented"):
    """Refuse --deliverables until the quality gates have actually run.

    The gates that guard a deck's quality — the design plan, the independent critic, the
    primary-source spot-check — are all *self-reported prose artifacts*. The model that skipped
    them is the same model that writes the hand-off note claiming they happened, so skipping is
    invisible: the note reads identically either way. This skill's own enforcement invariant
    ranks prose last for exactly that reason (deterministic lint > required-field > checklist >
    prose), yet these three sat at the bottom of it.

    So put the check where the incentive is. `--deliverables` produces the PDF and the preview
    page — the artifacts the user actually asked for — which makes it the one step nobody skips.
    Requiring evidence *here* costs nothing on a deck that ran its loop and blocks the shortcut on
    one that didn't.

    Evidence lives in `<deck-dir>/.deck-gates.json`:

        {"critic": {"verdict": "consent", "rounds": 2, "review": "reviews/r2.json"},
         "provenance": {"checked": 87, "confirmed": 85, "fixed": 2, "cut": 0}}

    A gate may be waived — quick decks exist — but a waiver is a written reason that travels with
    the deck, not a silence:

        {"critic": {"waived": "1-slide fix to an already-reviewed deck"}}

    Set SLIDE_MAKER_SKIP_GATES=1 to bypass entirely (CI smoke tests, throwaway renders).
    """
    if os.environ.get("SLIDE_MAKER_SKIP_GATES"):
        return
    path = os.path.join(os.path.dirname(os.path.abspath(pptx)) or ".", GATES_FILE)
    if not os.path.isfile(path):
        die("--deliverables is the hand-off run, and this deck has no record that its quality "
            "gates ran.\n"
            "  Missing: {}\n\n"
            "  The independent critic (Step 5) is the one that cannot be self-certified — you are "
            "not the final judge of your own deck.\n"
            "  Let the loop WRITE it, so the record is evidence and not a claim — one flag on a\n"
            "  step you already run before acting on any review:\n\n"
            "    python3 scripts/validate_review.py critic <review.json> --record {}\n\n"
            "  (that fills `critic` from the review itself, with its path + sha256, and the gate\n"
            "  re-reads it). Then add the two blocks only you can supply. Full shape:\n\n"
            '    {{"critic": {{"verdict": "consent", "rounds": 2}},\n'
            '     "design_plan": {{"boldness": "balanced+", "signature_move": "<the one risk>",\n'
            '                     "carried_by": [4, 6, 8], "form_ledger": "<family tally>",\n'
            '                     "icon_family": "<family | none - reason>",\n'
            '                     "palette": "<FILL vs TEXT-safe split, per palette_audit.py>",\n'
            '                     "type_scale": {{"display": 34, "title": 24, "body": 14}},\n'
            '                     "signature_proof": {{"slide": 6, "png": "render/slide06.png"}}}},\n'
            '     "provenance": {{"claims": [{{"claim": "<the claim>", "verdict": "CONFIRMED",\n'
            '                                "url": "https://<primary source>"}}]}}}}\n\n'
            "  (A summary tally like {{\"checked\": 87}} is REJECTED on purpose — a tally is "
            "written by the same pass that would have skipped the refutation.)\n\n"
            "  A gate you deliberately skipped is waived in writing, not omitted:\n\n"
            '    {{"critic": {{"waived": "why this deck does not need a critic round"}}}}\n\n'
            "  Renders without --deliverables are unaffected; iterate freely."
            .format(path, os.path.dirname(path) or "."))
    try:
        gates = json.load(open(path, encoding="utf-8"))
    except Exception as e:
        die("{} is not readable JSON: {}".format(path, e))

    critic = gates.get("critic") or {}
    if critic.get("waived"):
        # A waiver is legitimate — quick decks and hosts without subagent dispatch are real —
        # but an UNCLASSIFIED one is just a sentence, and the model that skipped the loop writes
        # the same sentence as the model that ran it. The Codex delivery gate has required a
        # distinct schema-valid review artifact per lens for a while; this path accepted any
        # string. Measured: a hand-typed waiver carried a whole deck through `all hand-off gates
        # pass` without an independent critic ever seeing it. So: name the category, and say
        # whether the lenses ran at all.
        WAIVER_KINDS = {
            "no-dispatch-on-host":
                "the runtime cannot dispatch a subagent (record it in the capability ledger too)",
            "already-reviewed-minor-edit":
                "a 1-2 slide edit to a deck that already passed its loop",
            "user-waived":
                "the user was asked and chose to ship over it",
            "external-deck":
                "a deck this skill did not author (redesign diagnosis / critique-only run)",
        }
        reason = critic["waived"]
        kind = critic.get("waived_category")
        if not isinstance(reason, str) or len(reason.strip()) < 24:
            die("`critic.waived` must be a written reason that travels with the deck, not a "
                "placeholder. Say what was skipped and why, in a sentence someone can disagree "
                "with later.")
        if kind not in WAIVER_KINDS:
            die("`critic.waived_category` must name WHICH kind of skip this is — an unclassified "
                "waiver is indistinguishable from never having run the loop. One of:\n"
                + "\n".join("    {:28s} {}".format(k, v) for k, v in sorted(WAIVER_KINDS.items()))
                + "\n\n  If none of these fit, the honest move is to run the critic.")
        if kind == "no-dispatch-on-host" and "inline_ran" not in critic:
            die("`waived_category: no-dispatch-on-host` must also record `\"inline_ran\": true|false` "
                "— whether the content and design lenses were at least run inline. 'Ran inline in "
                "the author's own context' and 'was never reviewed' are different claims, and the "
                "hand-off note reads identically for both unless this file separates them.")
        print("[gates] critic WAIVED [{}] — NOT INDEPENDENTLY REVIEWED".format(kind))
        print("        {}".format(reason))
        if kind == "no-dispatch-on-host":
            print("        lenses run inline: {} — inline review is the author grading "
                  "themselves.".format("yes" if critic.get("inline_ran") else "NO"))
        print("        Say this in the hand-off note too; a waiver the user never sees is a "
              "silence.")
    elif critic.get("verdict") == "consent":
        # A record the model TYPED at hand-off is self-certification: the model that skipped the
        # loop writes the same JSON as the model that ran it. So when the record points at the
        # review artifact (validate_review.py --record puts it there), re-read that artifact and
        # verify it still says what the record claims. A hand-written record still passes — but it
        # is LABELLED self-reported, so the two are distinguishable instead of identical.
        src = critic.get("source")
        if src:
            if not os.path.isfile(src):
                die("`critic.source` points at {} — which does not exist. The gate re-reads the "
                    "review artifact rather than trusting the summary; restore the file, re-run "
                    "`validate_review.py critic <review.json> --record <deck-dir>`, or waive in "
                    "writing.".format(src))
            digest = critic.get("sha256")
            if digest and _sha256(src) != digest:
                die("`critic.source` ({}) has changed since it was recorded — the sha256 no longer "
                    "matches. Re-validate it with `validate_review.py critic <review.json> "
                    "--record <deck-dir>` so the record and the evidence agree.".format(src))
            try:
                review = json.load(open(src, encoding="utf-8"))
            except Exception as exc:
                die("`critic.source` ({}) is not readable JSON: {}".format(src, exc))
            if review.get("verdict") != "consent":
                die("the recorded verdict says consent, but the review at {} says {!r}. The "
                    "artifact wins.".format(src, review.get("verdict")))
            hard = [f for f in (review.get("findings") or [])
                    if isinstance(f, dict) and f.get("severity") in ("blocker", "major")]
            if hard:
                die("the review at {} consents while still carrying {} blocker/major finding(s) — "
                    "that is a contract violation (agents/critic.md: any blocker/major -> revise). "
                    "Fix them and re-review, or waive in writing.".format(src, len(hard)))
            print("[gates] critic consented after {} round(s) — verified against {}".format(
                critic.get("rounds", "?"), os.path.basename(src)))
        else:
            print("[gates] critic consented after {} round(s) — SELF-REPORTED (no review "
                  "artifact).\n"
                  "        The evidence-backed path is one flag on a step you already run:\n"
                  "          python3 scripts/validate_review.py critic <review.json> --record {}"
                  .format(critic.get("rounds", "?"), os.path.dirname(path) or "."))
        if critic.get("corroborated_by"):
            print("[gates] consent corroborated by {} arbiter pass(es)".format(
                len(critic["corroborated_by"])))
    elif critic.get("verdict") == "revise":
        die("the last critic review returned verdict=revise. Fix the blockers and re-run the "
            "loop, or record a waiver with the reason you are shipping over it:\n"
            '    {"critic": {"waived": "shipping over: <the surviving finding and why>"}}')
    else:
        die("{} has no usable `critic` record — needs {{\"verdict\": \"consent\"|\"revise\"}} or "
            "{{\"waived\": \"<reason>\"}}".format(path))

    # The design plan is the art director's output (Step 2). Self-authoring one is indistinguishable
    # from dispatching for it — unless the record has to carry the fields the dispatch produces.
    # `icon_family` joins the four originals for one measured reason: a deck shipped with ZERO
    # icons through every automated gate, while a missing LOGO was caught by a required
    # checkpoint token. Icons are called a design must on every branch and had no field, no
    # column and no check anywhere. The token grammar mirrors `logo plan:` — a family name or
    # an explicit `none — <reason>` — so a deliberately icon-free deck is always satisfiable.
    # `palette` joins them for the same measured reason as `icon_family`. SKILL.md already
    # states the two-token rule -- a hue used as TEXT must itself clear 4.5:1, so keep a bright
    # FILL token and a darker TEXT twin. The rule is correct and still easy to break, because
    # the check is per-PAIR and a build touches dozens of them. Measured on one deck: the author
    # declared the rule in the design plan and then broke it FOUR times -- a vivid ochre set as a
    # label on its own pale slab (2.34:1), coral emphasis text on a coral tint (4.19:1), a table
    # highlight (4.19:1), a muted grey carrying real content on cream (4.26:1). None were
    # reckless; each was a pair nobody was thinking about while computing contrast for a
    # different pair, and each surfaced at render time or in review, a round later.
    # `scripts/palette_audit.py` resolves the whole matrix in one call, so the field is cheap to
    # fill honestly and cannot be filled at all without having run something.
    # type_scale and signature_proof were gated on the CODEX delivery path only, so on the shared
    # path typography was the one pillar of the visual language nobody had to resolve (palette,
    # icons and forms all had required fields), and the signature move was accepted as a SENTENCE
    # with nothing showing it survived into the render. That asymmetry is the exact shape of a bug
    # this repo already fixed once: the critic waiver was schema-checked for Codex and a hand-typed
    # string everywhere else, and it carried a real deck through "all gates pass".
    DESIGN_FIELDS = ("boldness", "signature_move", "carried_by", "form_ledger", "icon_family",
                     "palette", "type_scale", "signature_proof")
    # THE RESTRAINT CARVE — built on the escape the skill ALREADY documents, not a new one.
    # agents/slide-design.md: under a *conservative* dial (user-requested or purpose-defaulted for a
    # sober defense / regulatory / status deck) "the risk is OPTIONAL: take a modest, restrained
    # signature move if one fits, OR — if none does — fill the field with the one-clause
    # `deliberately restrained: <why>` so the field is never blank either way."
    #
    # That escape existed in prose and nowhere in the gate, so an honest 5-minute lab-meeting plan
    # was rejected for lacking a rendered `signature_proof` — a PNG proving an aesthetic risk
    # survived the build, for a deck that deliberately took none. The only way out was
    # {"waived": …}, which also switches off palette, type_scale and icon_family: the real choice
    # was "invent a risk" or "abandon all design gating", and both corrupt the record.
    #
    # What it deliberately does NOT do: it does not relax `signature_move` (the field is still never
    # blank — you must WRITE why restraint is the position) and it does not relax `carried_by`,
    # `palette`, `type_scale`, `icon_family` or `form_ledger`. And it cannot be claimed above the
    # conservative dial: at balanced+ and above a real signature move is required, not optional.
    design = gates.get("design_plan") or {}
    _dial = str(design.get("boldness", "")).strip().lower()
    _move = str(design.get("signature_move", "")).strip().lower()
    _carved = _dial == "conservative" and _move.startswith("deliberately restrained")
    if design.get("waived"):
        print("[gates] design plan WAIVED — {}".format(design["waived"]))
    elif design:
        required = [f for f in DESIGN_FIELDS if not (_carved and f == "signature_proof")]
        missing = [f for f in required if not design.get(f)]
        if missing:
            hint = ("\n    palette: the resolved FILL-only vs TEXT-safe split — run\n"
                    "      python3 scripts/palette_audit.py --from-style <deck>/style.py\n"
                    "    and paste what it hands back. A hue that works as a fill can read at\n"
                    "    2-4:1 as text on the same tint; the matrix is what stops that being\n"
                    "    found a render later." if "palette" in missing else "")
            die("`design_plan` is missing {}. These are the art director's outputs "
                "(agents/slide-design.md, Step 2) — a plan without them was not designed, it was "
                "defaulted. Fill them, or waive with a reason.{}".format(", ".join(missing), hint))
        scale = design["type_scale"]
        if not isinstance(scale, dict) or not all(
                isinstance(scale.get(k), (int, float)) for k in ("display", "title", "body")):
            die('`type_scale` must resolve the three tiers as numbers, e.g. '
                '{"display": 34, "title": 24, "body": 14}. SIZE SPRAWL tells authors to draw sizes '
                '"from the deck\'s declared type-scale tokens" — this is where they get declared. '
                'A deck with no scale does not have restrained typography, it has whatever each '
                'slide happened to pick.')
        # same floors the Codex delivery gate uses, so the two paths cannot disagree about what
        # counts as legible body type
        BODY_FLOORS = {"presented": 13.5, "textheavy": 13.5, "selfread": 12.0}
        delivery = str(gates.get("delivery", "presented"))
        floor = BODY_FLOORS.get(delivery, BODY_FLOORS["presented"])
        if scale["body"] < floor:
            die(f'`type_scale.body` is {scale["body"]}pt, under the {floor}pt floor for a '
                f'{delivery} deck — that is a legibility floor, not a style choice.')
        if not (scale["display"] > scale["title"] > scale["body"]):
            die(f'`type_scale` is not a scale: display {scale["display"]} > title {scale["title"]} '
                f'> body {scale["body"]} must hold, or the tiers do not rank and the hierarchy is '
                f'decorative rather than structural.')
        if _carved:
            print("[gates] boldness=conservative with a 'deliberately restrained:' signature move — "
                  "signature_proof not required (there is no risk to prove); every other field is")
        proof = None if _carved else design["signature_proof"]
        if proof is not None:
            # `png` or `path`: the Codex delivery gate spells this key `path` in its own evidence file,
            # and a Codex run keeps BOTH records (references/codex-runtime.md). Demanding one spelling
            # here would reject the field an OpenAI-bridged run naturally writes — the same evidence,
            # rejected for its key name.
            proof_file = (proof or {}).get("png") or (proof or {}).get("path")
            if not isinstance(proof, dict) or not proof_file or not proof.get("slide"):
                die('`signature_proof` must be {"slide": <n>, "png": "<rendered png>"} (the key may also '
                    'be spelled "path", as the Codex delivery gate does) — the rendered evidence that '
                    'the signature move actually SURVIVED into the deck. A move that exists only as a '
                    'sentence in the plan is the documented failure: it gets sanded back to the safe '
                    'catalogue during the build and nobody notices, because the plan still reads bravely.')
            png = Path(proof_file)
            if not png.is_absolute():
                png = Path(pptx).parent / png
            if not png.exists() or png.stat().st_size < 512:
                die(f'`signature_proof` file ({proof_file}) does not exist or is empty — render '
                    f'the slide first (render_deck.py <deck> <dir>) and point at the real PNG.')
        cb = design["carried_by"]
        if not isinstance(cb, list) or len(cb) < 2:
            die("`carried_by` must name at least 2 slides where the signature move does structural "
                "work. One brave slide among eleven safe ones is a tonal break, not a position.")
        print("[gates] design plan: boldness={} · signature={} · carried_by={}".format(
            design["boldness"], str(design["signature_move"])[:48], cb))
    else:
        die("no `design_plan` record. Step 2 dispatches agents/slide-design.md as the deck's art "
            "director; nothing else decides deck rhythm or the signature move.\n"
            '    {"design_plan": {"boldness": "balanced+", "signature_move": "...",\n'
            '                     "carried_by": [4, 6, 8], "form_ledger": "...",\n'
            '                     "icon_family": "...", "palette": "..."}}\n'
            '    or {"design_plan": {"waived": "<reason>"}}')

    # Provenance: a self-filled tally proves nothing — the refutation pass is what the gate is FOR.
    # Require per-claim verdicts, so "confirmed" means someone tried to break it and could not.
    prov = gates.get("provenance") or {}
    if prov.get("waived"):
        print("[gates] provenance WAIVED — {}".format(prov["waived"]))
    elif prov:
        claims = prov.get("claims")
        if not isinstance(claims, list) or not claims:
            die("`provenance` needs a per-claim `claims` list, not a summary tally. A tally is "
                "written by the same pass that would have skipped the check.\n"
                '    "claims": [{"claim": "...", "verdict": "CONFIRMED|WRONG|PARTLY-WRONG|'
                'UNVERIFIABLE", "url": "https://..."}]')
        ALLOWED = {"CONFIRMED", "WRONG", "PARTLY-WRONG", "UNVERIFIABLE"}
        bad = [c for c in claims if c.get("verdict") not in ALLOWED]
        if bad:
            die("{} claim row(s) carry no valid verdict (one of {}).".format(
                len(bad), " / ".join(sorted(ALLOWED))))
        nourl = [c for c in claims if c.get("verdict") == "CONFIRMED" and not c.get("url")]
        if nourl:
            die("{} claim(s) are CONFIRMED with no primary-source URL. Confirmed-without-a-source "
                "is the exact failure this gate exists to catch.".format(len(nourl)))
        unresolved = [c for c in claims if c["verdict"] in ("WRONG", "PARTLY-WRONG")]
        if unresolved:
            die("{} claim(s) still verdict WRONG / PARTLY-WRONG. Fix or cut them before hand-off:\n"
                "  - {}".format(len(unresolved),
                                "\n  - ".join(str(c.get("claim"))[:90] for c in unresolved[:5])))
        tally = {}
        for c in claims:
            tally[c["verdict"]] = tally.get(c["verdict"], 0) + 1
        print("[gates] provenance: {} claim(s) adversarially checked — {}".format(
            len(claims), " · ".join("{} {}".format(v, k) for k, v in sorted(tally.items()))))
    else:
        print("[gates] no provenance record — fine for a deck built from the user's own material; "
              "a research-sourced deck should carry one.")

    # ── DENSITY: a slide is a visual aid, not a document ────────────────────────
    # This one is a gate rather than a warning because the warning already existed and was
    # already ignored — twice, by the same author, on two consecutive decks. Measured: one
    # deck shipped with 8 of 12 slides over the presented text budget, the next with 12 of 12
    # (loads of 81-144 words against a budget of ~40), and both times the per-slide TEXT WALL
    # line was read and dismissed as advisory. The skill's OWN reference deck — the file it
    # tells every builder to copy — runs at a median of 27 words a slide, so the budget is not
    # unrealistic; what failed was that nothing made ignoring it cost anything.
    # A deck may legitimately be denser (a self-read leave-behind, a spec sheet). That is what
    # the waiver is for: not a rule against text, a rule against text arriving by default.
    #
    # DELIVERY MODE. The budget here is lint_deck's budget, taken from lint_deck: 70 words for a
    # presented deck, 120 for a self-read one, and no budget at all on a poster (`surface`) or a
    # deck whose density the user CHOSE at Q4 (`textheavy`). A gate that fires on a mode the
    # interview offers, the rubric protects and the lint deliberately passes does not enforce
    # anything for long — it teaches the author to paste a waiver, and after that it is decoration.
    txt = gates.get("density")
    if mode in ("surface", "textheavy"):
        print("[gates] density: not applied — %s deck (the user chose this density, or the "
              "surface has no per-slide budget)" % mode)
        return
    over, total, median = _density_stats(pptx, budget=70 if mode == "presented" else 120)
    if total:
        if isinstance(txt, dict) and txt.get("waived"):
            print("[gates] density: {}/{} slide(s) over the presented text budget, median {} "
                  "words — WAIVED: {}".format(over, total, median, str(txt["waived"])[:110]))
        elif over * 3 > total:
            die("{} of {} slides are over the {} text budget (median {} words a slide; "
                "aim ~40, warn >{}). The skill's own reference deck runs at 27.\n"
                "    A slide is a visual aid for a speaker — the sentences belong in the speaker "
                "notes, which this deck already has.\n"
                "    Cut the on-slide prose, or record the deliberate choice:\n"
                '    "density": {{"waived": "why this deck is meant to be read, not presented"}}'
                .format(over, total, mode, median, 70 if mode == "presented" else 120))
        else:
            print("[gates] density: {}/{} slide(s) over the text budget, median {} words a slide"
                  .format(over, total, median))


def _density_stats(pptx, budget=70):
    """(slides over the text budget, total slides, median load).

    Calls `lint_deck.reading_load` — the SAME function the per-slide TEXT WALL warning calls —
    so the number in the warning and the number in this gate are one number by construction.
    An earlier version of this only shared the string-level word counter and re-implemented the
    per-slide accumulation and the chrome filter, which is where all the drift actually lives:
    its `sz <= 10.5 and len(t) < 40` skip had no POSITION test, so it was not a footer filter
    but an amnesty for small type anywhere on the slide. Measured on one deck: lint said 136
    words a slide, the gate said 4, and a wall of 10.5pt prose passed. Sharing a helper is not
    the same as sharing the measurement — share the measurement.
    """
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from lint_deck import _boxes, reading_load
        from pptx import Presentation
        prs = Presentation(pptx)
        sw = prs.slide_width / 914400.0
        sh = prs.slide_height / 914400.0
    except Exception:
        return 0, 0, 0
    loads = [reading_load(sl, _boxes(sl, sw, sh), sh) for sl in prs.slides]
    if not loads:
        return 0, 0, 0
    loads.sort()
    return sum(1 for x in loads if x > budget), len(loads), loads[len(loads) // 2]


def main(argv):
    # --deliverables (alias --final): ALSO park the PDF beside the .pptx and write viewer.html.
    # OFF by default: while a deck is still being iterated, those two are pure churn — they are
    # regenerated every round, clutter the deck root, and go stale the moment the user hand-edits
    # the .pptx. They are produced once, at hand-off, when the user says the deck is final.
    deliverables = False
    # --fast: re-render ONLY the slides whose content changed since the last render.
    # Measured on a real 18-slide deck: full = 9.1s LibreOffice + 24.4s rasterize; one changed
    # slide = 2.5s + 0.7s. Rasterization, not the PDF export, is the dominant cost — so the win
    # comes from rasterizing one page, and subsetting the pptx makes the export cheap too.
    fast = False
    # --slides N[,M]: render ONLY these 1-indexed slides. Unlike --fast (which DIFFS against a
    # cache and therefore renders everything on a first run), this is an explicit "show me page N"
    # — the SIGNATURE PROOF that opens SKILL.md Step 4, and any "re-render the page I edited" loop.
    only = None
    argv = list(argv)
    while "--fast" in argv:
        argv.remove("--fast")
        fast = True
    for i, a in enumerate(list(argv)):
        if a == "--slides" or a.startswith("--slides="):
            raw = a.split("=", 1)[1] if "=" in a else (argv[i + 1] if i + 1 < len(argv) else "")
            try:
                only = sorted({int(t) for t in raw.replace(" ", "").split(",") if t})
            except ValueError:
                die("--slides wants 1-indexed slide numbers, e.g. --slides 1,6")
            if not only:
                die("--slides wants at least one slide number, e.g. --slides 6")
            argv.remove(a)
            if "=" not in a and raw in argv:
                argv.remove(raw)
            break
    # --gate-check runs the hand-off gates and exits, rendering nothing. The gates were reachable
    # only through --deliverables, which Step 6 deliberately makes a decline-able OFFER ("want a PDF
    # and a browser preview?"), so on every deck where the user said no, the strongest gate in the
    # skill never ran at all. Step 6 now calls this unconditionally, whatever the user answers.
    # Delivery mode — the SAME flags lint_deck.py takes, spelled the same way, because the gate
    # below enforces the lint's budget and the two must never disagree about which deck this is.
    mode = "presented"
    for flag, name in (("--selfread", "selfread"), ("--surface", "surface"),
                       ("--textheavy", "textheavy")):
        while flag in argv or ("--mode=" + name) in argv:
            argv = [a for a in argv if a not in (flag, "--mode=" + name)]
            mode = name
    gate_only = False
    if "--gate-check" in argv:
        argv = [a for a in argv if a != "--gate-check"]
        gate_only = True
    for flag in ("--deliverables", "--final"):
        while flag in argv:
            argv.remove(flag)
            deliverables = True
    if fast and deliverables:
        # Decided here, not after LibreOffice has already run: --deliverables needs a whole-deck
        # PDF, which a subset render cannot produce. Failing late meant either an exit-1 after a
        # successful render, or (with nothing changed) a silent exit-0 that produced no PDF and no
        # viewer.html at the exact moment the hand-off contract required them.
        die("--deliverables needs a full-deck render — drop --fast for the hand-off run")
    if only and deliverables:
        # Same reason as --fast: a subset cannot produce the whole-deck PDF the hand-off promises.
        die("--deliverables needs a full-deck render — drop --slides for the hand-off run")
    if only and fast:
        # Contradictory intents: --fast decides WHICH slides to render, --slides declares them.
        die("--slides and --fast both choose the slide set — pass one")
    if not argv:
        die("usage: python render_deck.py /path/to/deck.pptx [out_dir] "
            "[--fast | --slides N[,M]] [--deliverables] [--gate-check]")
    pptx = argv[0]
    out = argv[1] if len(argv) > 1 else "./render"

    if not os.path.isfile(pptx):
        die("no such file: " + pptx)

    if gate_only:
        check_handoff_gates(pptx, mode)
        print("[gates] all hand-off gates pass — the deck may be handed over")
        return 0

    # Checked here, before LibreOffice runs — same reason the --fast/--slides conflicts are:
    # failing after a successful render wastes the render and reads as a late surprise.
    if deliverables:
        check_handoff_gates(pptx, mode)

    soffice = find_soffice()
    if not soffice:
        die(
            "LibreOffice not found — needed to render slides for the verify + critic loop.\n"
            "  macOS:         brew install --cask libreoffice\n"
            "  Debian/Ubuntu: sudo apt install libreoffice\n"
            "  Windows:       winget install TheDocumentFoundation.LibreOffice\n"
            "                 (or choco install libreoffice-fresh)\n"
            "  other:         https://www.libreoffice.org/download\n"
            "  (or set the SOFFICE env var to the full path of the soffice binary)"
        )

    # Decide full vs incremental BEFORE spending anything on LibreOffice.
    # When out == the deck's own folder (the code below explicitly tolerates the user passing "."),
    # <deck>.pdf and viewer.html there are the HAND-OFF deliverables, not render products. Nothing
    # in this run may delete or overwrite them.
    _deck_dir = os.path.dirname(os.path.abspath(pptx)) or "."
    _out_is_deck_dir = os.path.abspath(out) == os.path.abspath(_deck_dir)

    cache_path = os.path.join(out, ".render-cache.json")
    # Cheap, but not free on a big deck — and a full render still needs them, both to validate the
    # page count and to leave a cache the NEXT --fast can diff against. Guarded so an unreadable
    # package falls back to LibreOffice's own diagnostic instead of a zipfile traceback (v3.5.1
    # behaviour, which the friendly error in troubleshooting-faq.md documents).
    try:
        fps, blockers = _slide_fingerprints(pptx)
    except Exception as exc:
        fps, blockers = [], ["could not read the .pptx package: {}".format(exc)]
    changed, skip_reason = None, None
    if fast:
        prev = None
        try:
            with open(cache_path, encoding="utf-8") as f:
                cached = json.load(f)
            prev = cached.get("fingerprints") if isinstance(cached, dict) else None
            if not (isinstance(prev, list) and all(isinstance(x, str) for x in prev)):
                prev = None                     # a JSON list/garbage must not crash the tool
        except Exception:
            prev = None
        if blockers:
            # Correctness beats speed: anything that makes the slide->page mapping uncertain
            # forces a full render.
            skip_reason = blockers[0]
        elif not prev:
            skip_reason = "no previous render cache"
        elif len(prev) != len(fps):
            # Slides inserted or deleted shifts every index after the edit; PNG filenames would
            # silently point at the wrong slides.
            skip_reason = "slide count changed ({} -> {})".format(len(prev), len(fps))
        else:
            changed = [i for i, h in enumerate(fps) if h != prev[i]]

            def _usable(name):                  # a 0-byte file from a killed rasterize is NOT a render
                p = os.path.join(out, name)
                try:
                    return os.path.getsize(p) > 0
                except OSError:
                    return False

            missing = [i for i in range(len(fps)) if not _usable("slide{:02d}.png".format(i + 1))]
            # thumbnails are regenerated only when a bookend slide is re-rendered, so a missing
            # thumb must pull its bookend into the changed set or the critic's poster test runs
            # on a stale or absent image while the run prints success
            if not _usable("thumb_first.png"):
                missing.append(0)
            if not _usable("thumb_last.png"):
                missing.append(len(fps) - 1)
            changed = sorted(set(changed) | set(i for i in missing if 0 <= i < len(fps)))

    try:
        _st = os.stat(pptx)
        pptx_stat = (_st.st_mtime_ns, _st.st_size)
    except OSError:
        pptx_stat = None
    if only is not None:
        if blockers:
            # The SAME correctness bar as --fast: anything that makes the slide->page mapping
            # uncertain (hidden slides, an unresolvable part, auto slide-number fields that
            # renumber inside a subset) must not be papered over just because the user named a page.
            die("cannot render a subset of this deck: {}\n"
                "  drop --slides and render the whole deck".format(blockers[0]))
        bad = [n for n in only if not (1 <= n <= len(fps))]
        if bad:
            die("--slides {} out of range — this deck has {} slide(s)".format(
                ",".join(str(b) for b in bad), len(fps)))
        changed = [n - 1 for n in only]
    incremental = (fast or only is not None) and changed is not None and 0 < len(changed) < len(fps)
    if fast and changed is not None and len(changed) == len(fps) and len(fps):
        # Every slide changed (a rebuild that touched everything, or a deck-global edit such as a
        # theme/canvas change). A subset of "all slides" is just a full render with extra steps.
        skip_reason = "every slide changed"
    if fast and changed is not None and not changed:
        print("no slide changed since the last render — nothing to re-render")
        print("next: python3 {} {} --renders {}".format(
            os.path.join(os.path.dirname(os.path.abspath(__file__)), "lint_deck.py"), pptx, out))
        return 0



    # Give this invocation its OWN LibreOffice profile: lets parallel renders (the
    # large-deck section fan-out) run at once without fighting a shared profile lock,
    # and lets the render work even while the user has the LibreOffice GUI open.
    # Without this, concurrent/coexisting soffice calls silently produce no PDF.
    src_pptx, keep = pptx, None
    tmp_subset = tmp_dir = None
    if incremental:
        keep = changed
        tmp_dir = tempfile.mkdtemp(prefix="lo_subset_")
        tmp_subset = os.path.join(tmp_dir, "subset.pptx")
        _subset_pptx(pptx, set(keep), tmp_subset)
        src_pptx = tmp_subset

    # A subset renders into its OWN temp dir. Writing out/subset.pdf would collide between two
    # concurrent runs sharing a render dir (which the per-invocation LibreOffice profile above
    # exists to allow), and a crashed run would leave a file that breaks the render-only cleanup.
    # ALWAYS convert into a private empty directory, never straight into `out`: `out` may already
    # hold a <deck>.pdf (a previous render, or a --deliverables artefact when out is the deck
    # folder), and a stale file there would make a failed conversion look like a successful one.
    if tmp_dir is None:
        tmp_dir = tempfile.mkdtemp(prefix="lo_render_out_")
    import atexit
    atexit.register(shutil.rmtree, tmp_dir, True)       # a die() mid-render must not leak a deck copy
    pdf, result, cmd = _render_pdf(soffice, src_pptx, tmp_dir)
    try:
        import fitz  # pymupdf
    except ImportError:
        die("pymupdf not installed — run: {} -m pip install pymupdf".format(
            os.path.basename(sys.executable) or "python"))

    if result.returncode != 0 or not os.path.isfile(pdf):
        detail = [
            "LibreOffice failed to convert {} (exit {}).".format(pptx, result.returncode)
            if result.returncode != 0 else "LibreOffice produced no PDF from {}.".format(pptx),
            "Command: " + " ".join(cmd),
            "Exit code: {}".format(result.returncode),
        ]
        stdout = _tail(result.stdout)
        stderr = _tail(result.stderr)
        if stdout:
            detail.append("stdout:\n" + stdout)
        if stderr:
            detail.append("stderr:\n" + stderr)
        detail.append(
            "Check that the file opens, close any open copy, and in sandboxed runtimes "
            "rerun the render with the permissions needed for LibreOffice."
        )
        die("\n".join(detail))

    # Open the PDF BEFORE the cleanup: a non-zero exit and a missing file were already fatal, but a
    # zero-exit run that writes a truncated/garbage PDF passed both checks, and the cleanup had
    # already deleted the previous render by the time fitz raised.
    try:
        _probe = fitz.open(pdf)
        _npages = _probe.page_count
        _probe.close()
    except Exception as exc:
        die("LibreOffice produced an unreadable PDF from {} ({}). The previous render in {} was "
            "left untouched.".format(pptx, exc, out))
    if _npages < 1:
        die("LibreOffice produced an empty PDF (0 pages) from {}. The previous render in {} was "
            "left untouched.".format(pptx, out))

    # Clear the previous render only NOW, after LibreOffice has actually produced a readable PDF. Doing it
    # earlier meant a failed conversion destroyed the last good render and left nothing at all —
    # the user lost working output to a run that produced none.
    # Wiping the render dir BEFORE deciding is what made --fast a no-op: every PNG looked
    # "missing", so every slide looked changed. An incremental run keeps the existing PNGs and
    # overwrites only the ones it re-renders.
    if os.path.isdir(out) and not incremental:
        entries = os.listdir(out)
        own_pdf = os.path.splitext(os.path.basename(pptx))[0] + ".pdf"
        render_only = all(
            (e.startswith(("slide", "thumb_")) and e.endswith(".png"))
            or e == "viewer.html" or e == own_pdf          # only THIS deck's fallback pdf is ours
            or e in (".DS_Store", "Thumbs.db", ".render-cache.json")   # OS junk + our own cache only
            for e in entries)                              # make the dir look deletable
        if render_only:
            shutil.rmtree(out, ignore_errors=True)
        else:
            # out holds files that are NOT ours (worst case: the user passed "." — the pptx's own
            # directory). NEVER rmtree it; clear only our previous render products.
            for e in entries:
                if _out_is_deck_dir and (e == "viewer.html" or not _RENDER_PNG.match(e)):
                    # out IS the deck folder: only files matching our own strict output pattern
                    # (slideNN.png / thumb_first|last.png) are ours. A user's slide_background.png
                    # or thumb_hero.png sitting beside the deck was being deleted silently.
                    continue
                if (e.startswith(("slide", "thumb_")) and e.endswith(("png",))) or e == "viewer.html":
                    try:
                        os.remove(os.path.join(out, e))
                    except OSError:
                        pass
    os.makedirs(out, exist_ok=True)


    doc = fitz.open(pdf)
    skip_cache = False
    if only is not None:
        # A --slides run rasterized SOME pages; the fingerprints in hand describe ALL of them.
        # Caching them would tell the next --fast "nothing changed" while most PNGs are stale —
        # precisely the lie the incremental path exists to avoid. Drop the cache instead, so the
        # next --fast honestly does a full render.
        skip_cache = True
    if incremental:
        # PDF page k corresponds to deck slide keep[k] — write ONLY those PNGs and leave the
        # rest of render/ untouched.
        if doc.page_count != len(keep):
            die("incremental render produced {} page(s) for {} changed slide(s) — refusing to "
                "write PNGs that may be mismatched; re-run without --fast".format(
                    doc.page_count, len(keep)))
        for k, page in enumerate(doc):
            page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(
                os.path.join(out, "slide{:02d}.png".format(keep[k] + 1)))
        # thumbnails only matter when a bookend slide moved
        for name, idx in (("thumb_first", 0), ("thumb_last", len(fps) - 1)):
            if idx in keep:
                page = doc[keep.index(idx)]
                zoom = 240.0 / max(1.0, page.rect.width)
                page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(os.path.join(out, name + ".png"))
        n_pages = len(fps)
    else:
        pages = list(enumerate(doc, 1))
        for i, page in pages:
            page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(
                os.path.join(out, "slide{:02d}.png".format(i)))
        # bookend thumbnails (~240px wide) for the critic's poster test — first + last slide small,
        # the scale at which a cover either survives or dies. Same PyMuPDF path, no new deps.
        if pages:
            for name, (_, page) in (("thumb_first", pages[0]), ("thumb_last", pages[-1])):
                zoom = 240.0 / max(1.0, page.rect.width)
                page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(os.path.join(out, name + ".png"))
        n_pages = doc.page_count
        if n_pages != len(fps) and fps:
            # LibreOffice emitted a different number of pages than the deck has slides (hidden
            # slides are the known cause). slideNN.png no longer means deck slide NN, so no cache
            # may claim this render — and the user must be told rather than quietly trusting it.
            print("WARNING: {} slides in the deck but {} PDF page(s) rendered — slideNN.png may "
                  "not correspond to deck slide NN. Not caching fingerprints; --fast is disabled "
                  "until this is resolved.".format(len(fps), n_pages), file=sys.stderr)
            skip_cache = True
    doc.close()
    if incremental:
        pdf = None                              # a subset PDF is not the deck's PDF
        # Drop the previous FULL render's intermediate PDF, which the new PNGs have outdated —
        # but never when `out` IS the deck folder, where that same filename is the user's
        # hand-off deliverable. Deleting it there silently destroyed a --deliverables artefact.
        if not _out_is_deck_dir:
            try:
                os.remove(os.path.join(out, os.path.splitext(os.path.basename(pptx))[0] + ".pdf"))
            except OSError:
                pass

    # Record fingerprints so the NEXT run can diff against them. Written only after the PNGs
    # actually landed, so a crashed render never leaves a cache claiming work that did not happen.
    # Two guards: (1) if the .pptx changed WHILE we rendered, the fingerprints we hold describe
    # state that was never rasterized — caching them would freeze that slide stale forever;
    # (2) a failed write must DELETE the cache, because a cache older than the PNGs is exactly the
    # "no slide changed" lie this whole path must not tell.
    if not skip_cache and fps:
        # A stat() beats re-hashing the whole package: we only need to know whether the file moved
        # under us during the ~10s render, and mtime+size answers that for a fraction of the cost.
        try:
            st = os.stat(pptx)
            moved = (st.st_mtime_ns, st.st_size) != pptx_stat
        except OSError:
            moved = True
        if moved:
            skip_cache = True
            print("note: the .pptx changed during the render — not caching fingerprints",
                  file=sys.stderr)
    if skip_cache:
        try:
            os.remove(cache_path)
        except OSError:
            pass
    else:
        try:
            tmp_cache = cache_path + ".tmp"
            with open(tmp_cache, "w", encoding="utf-8") as f:
                json.dump({"fingerprints": fps}, f)
            os.replace(tmp_cache, cache_path)
        except OSError:
            try:
                os.remove(cache_path)
            except OSError:
                pass
            print("note: could not write the render cache — the next --fast will do a full render",
                  file=sys.stderr)

    # The PDF is an INTERMEDIATE of this render (pptx -> PDF -> PNG), so it always exists. Whether
    # it is promoted to a deliverable beside the .pptx is the user's call at hand-off.
    pdf_dest = pdf
    if deliverables and pdf is None:
        die("--deliverables needs a full-deck render; re-run without --fast")
    if deliverables:
        pdf_dest = os.path.join(os.path.dirname(os.path.abspath(pptx)) or ".",
                                os.path.splitext(os.path.basename(pptx))[0] + ".pdf")
        try:
            if os.path.abspath(pdf_dest) != os.path.abspath(pdf):
                shutil.move(pdf, pdf_dest)   # move, not replace: the source is a temp dir that may
        except OSError:                      # sit on a different filesystem
            pdf_dest = pdf                   # couldn't move (odd mount/permissions)

    # Self-contained flip-through viewer — parked BESIDE the .pptx (deck root), same as the PDF, so
    # the user finds it without digging into render/. It references the PNGs through the render subdir
    # (relative to the viewer's own location), so a plain double-click works. One file:// link, any
    # browser, any OS (arrow keys / click / thumbnail strip). Zero dependencies, zero network.
    deck_dir = os.path.dirname(os.path.abspath(pptx)) or "."
    viewer = None
    if deliverables:
        rel = os.path.relpath(os.path.abspath(out), deck_dir)
        pref = "" if rel in (".", "") else rel.replace(os.sep, "/").rstrip("/") + "/"
        slides = [pref + "slide{:02d}.png".format(i) for i in range(1, n_pages + 1)]
        viewer = os.path.join(deck_dir, "viewer.html")
        try:
            with open(viewer, "w", encoding="utf-8") as f:
                f.write(_viewer_html(os.path.splitext(os.path.basename(pptx))[0], slides))
        except OSError:
            viewer = None
    # sweep a stale viewer.html left inside the render dir by an older build (it now lives at root)
    stale = os.path.join(out, "viewer.html")
    if viewer and os.path.abspath(stale) != os.path.abspath(viewer) and os.path.exists(stale):
        try: os.remove(stale)
        except OSError: pass
    if incremental:
        print("{}: {} of {} slides re-rendered ({}) -> {}".format(
            "slides render" if only is not None else "fast render",
            len(keep), len(fps), ", ".join(str(i + 1) for i in keep), out))
        # If a hand-off already produced the deck-root pair, they now lag the deck. Say so loudly:
        # a stale PDF someone opens and reviews is the failure this whole path is built to avoid.

    else:
        print("rendered {} slides -> {}".format(n_pages, out))
        if fast and skip_reason:
            print("(--fast fell back to a full render: {})".format(skip_reason))
    if not deliverables:
        print("pdf/viewer: not generated (deck still in progress) — at hand-off, once the user "
              "confirms the deck is final, re-run with --deliverables")
    else:
        print("pdf: {}".format(pdf_dest))
    if viewer:
        print("preview: {}  (open in a browser; arrow keys flip)".format(
            Path(viewer).resolve().as_uri()))
    # Any render that is NOT the hand-off run leaves an already-delivered pair behind the deck.
    # This must fire on the full path too: re-rendering a delivered deck is the likeliest way to
    # end up with a PDF someone opens and reviews after it stopped matching the .pptx.
    if not deliverables:
        _stale = [f for f in (os.path.splitext(os.path.basename(pptx))[0] + ".pdf", "viewer.html")
                  if os.path.isfile(os.path.join(_deck_dir, f))]
        if _stale:
            print("WARNING: {} at the deck root {} now STALE — re-run with --deliverables (and "
                  "without --fast) before handing the deck over".format(
                      " and ".join(_stale), "is" if len(_stale) == 1 else "are"), file=sys.stderr)
    print("next: python3 {} {} --renders {}  # render-time lint, then the actor-critic loop".format(
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "lint_deck.py"), pptx, out))


def _viewer_html(title, slides):
    """Single-file dark flip-through viewer: big slide + thumbnail rail + keyboard/click nav."""
    import json
    return """<!DOCTYPE html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title} — preview</title>
<style>
  :root {{ color-scheme: dark; }}
  * {{ margin: 0; box-sizing: border-box; }}
  body {{ background: #17191f; color: #cfd4de; font: 13px/1.4 system-ui, sans-serif;
         height: 100vh; display: flex; flex-direction: column; user-select: none; }}
  header {{ padding: 8px 14px; display: flex; align-items: center; gap: 12px; }}
  header b {{ color: #fff; font-weight: 600; }}
  #stage {{ flex: 1; display: flex; align-items: center; justify-content: center;
            min-height: 0; padding: 0 12px; cursor: pointer; }}
  #main {{ max-width: 100%; max-height: 100%; box-shadow: 0 6px 30px rgba(0,0,0,.5);
           border-radius: 4px; }}
  #rail {{ display: flex; gap: 6px; overflow-x: auto; padding: 10px 14px; flex: none; }}
  #rail img {{ height: 62px; border-radius: 3px; opacity: .45; cursor: pointer;
               border: 2px solid transparent; }}
  #rail img.on {{ opacity: 1; border-color: #5b8def; }}
  #num {{ margin-left: auto; font-variant-numeric: tabular-nums; color: #8a93a6; }}
  kbd {{ background:#2a2e38; border-radius:3px; padding:1px 5px; font-size:11px; color:#9aa3b2; }}
</style></head><body>
<header><b>{title}</b><span>&larr;/&rarr; or click to flip &nbsp;<kbd>F</kbd> fullscreen</span><span id="num"></span></header>
<div id="stage"><img id="main" alt="slide"></div>
<div id="rail"></div>
<script>
const S = {slides}; let i = 0;
const main = document.getElementById('main'), rail = document.getElementById('rail'),
      num = document.getElementById('num');
S.forEach((src, k) => {{ const t = document.createElement('img'); t.src = src; t.loading = 'lazy';
  t.onclick = () => go(k); rail.appendChild(t); }});
function go(k) {{ i = (k + S.length) % S.length; main.src = S[i];
  num.textContent = (i + 1) + ' / ' + S.length;
  [...rail.children].forEach((t, k2) => t.classList.toggle('on', k2 === i));
  rail.children[i].scrollIntoView({{ inline: 'center', block: 'nearest', behavior: 'smooth' }});
  if (i + 1 < S.length) (new Image()).src = S[i + 1]; }}
document.getElementById('stage').onclick = () => go(i + 1);
addEventListener('keydown', e => {{
  if (e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown') go(i + 1);
  else if (e.key === 'ArrowLeft' || e.key === 'PageUp') go(i - 1);
  else if (e.key === 'Home') go(0); else if (e.key === 'End') go(S.length - 1);
  else if (e.key.toLowerCase() === 'f') document.documentElement.requestFullscreen?.(); }});
go(0);
</script></body></html>
""".format(title=title.replace("&", "&amp;").replace("<", "&lt;"), slides=json.dumps(slides))


if __name__ == "__main__":
    main(sys.argv[1:])
