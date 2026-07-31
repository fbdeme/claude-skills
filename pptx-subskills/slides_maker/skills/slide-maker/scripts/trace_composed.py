#!/usr/bin/env python3
"""Split a deck's shipped lines into SOURCE-TRACEABLE and AUTHOR-COMPOSED, so the content
review can aim at the second set instead of re-reading all twelve pages.

WHAT THIS IS NOT. It does not detect fabrication, and the first version of it tried to. That
version was measured on a real deck and abandoned: of 21 shipped lines it flagged 13, including
faithful paraphrases ("同一段内容，四种设计语言" — 7 absent trigrams) while the ONE actually invented
mechanism ("对不上就标为待查") had 5. Precision ~8%. A check that fires on what it cannot know is one
people learn to ignore, so that idea is deliberately not here.

WHAT THE SAME MEASUREMENT DID SHOW is a clean and different signal. The lines with ZERO absent
n-grams were, without exception, the ones lifted near-verbatim from the source; every flagged line
was the author's own composition. The split is not true/false — it is **quoted vs written**. And
that is worth having, because on that deck all three of the worst content defects landed in the
written set:

  · an invented product mechanism, on the page whose argument was "never invent"
  · a process gate attributed to the wrong step
  · a required install line silently dropped

None were detectable by rule. All three were the author writing from memory of the source hours
after reading it — which is exactly what "composed, not quoted" identifies. So this tool does the
one honest thing available: it narrows where a human or a critic has to look, and says plainly
that it is triage and not a verdict.

    python3 scripts/trace_composed.py <deck.pptx> --source README.md,SKILL.md [--n 3] [--json out]

Feed the COMPOSED list into the content critic's brief. A traceable line still needs its number
checked; a composed line needs its whole claim checked.

Exit 0 always — this is triage, and triage that fails a build would be dishonest about what it
knows. It reports NOT CHECKED (exit 2) when it cannot read its inputs.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

CJK = r"[㐀-䶿一-鿿豈-﫿]"


def slide_lines(path: Path):
    """(slide_no, paragraph) for every non-empty text paragraph, in deck order."""
    from pptx import Presentation
    prs = Presentation(str(path))
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if t:
                    yield i, t


def cjk_ngrams(s: str, n: int) -> set[str]:
    """n-grams over CJK runs only. Latin tokens and digits are checked separately and exactly —
    an n-gram model on them produces noise without adding information."""
    out = set()
    for run in re.findall(f"{CJK}+", s):
        for i in range(len(run) - n + 1):
            out.add(run[i:i + n])
    return out


def latin_tokens(s: str) -> set[str]:
    """Identifiers, commands, versions, paths — the class where a mismatch is objectively wrong.
    These get an exact substring test, not a fuzzy one."""
    return {t for t in re.findall(r"[A-Za-z][A-Za-z0-9._/+-]{2,}", s)
            if not t.isdigit() and len(t) > 2}


def numbers(s: str) -> set[str]:
    return set(re.findall(r"\d+(?:\.\d+)?%?", s))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--source", required=True,
                    help="comma-separated source files (or dirs) the deck must trace to")
    ap.add_argument("--n", type=int, default=3, help="CJK n-gram size (default 3)")
    ap.add_argument("--json", help="write the split to this path")
    a = ap.parse_args()

    deck = Path(a.deck)
    if not deck.is_file():
        print(f"NOT CHECKED — no such deck: {deck}", file=sys.stderr)
        return 2

    corpus, read = "", []
    for raw in a.source.split(","):
        p = Path(raw.strip()).expanduser()
        files = sorted(p.rglob("*.md")) if p.is_dir() else [p]
        for f in files:
            try:
                corpus += f.read_text(encoding="utf-8", errors="ignore")
                read.append(str(f))
            except OSError as e:
                print(f"NOT CHECKED — cannot read {f}: {e}", file=sys.stderr)
                return 2
    if not corpus:
        print(f"NOT CHECKED — no source text read from {a.source}", file=sys.stderr)
        return 2

    try:
        lines = list(slide_lines(deck))
    except Exception as e:
        print(f"NOT CHECKED — could not open {deck}: {e}", file=sys.stderr)
        return 2

    quoted, composed, hard = [], [], []
    for n_slide, txt in lines:
        absent_ng = sorted(g for g in cjk_ngrams(txt, a.n) if g not in corpus)
        absent_tok = sorted(t for t in latin_tokens(txt) if t not in corpus)
        absent_num = sorted(v for v in numbers(txt) if v not in corpus)
        rec = {"slide": n_slide, "text": txt, "absent_ngrams": absent_ng,
               "absent_tokens": absent_tok, "absent_numbers": absent_num}
        # a Latin identifier or a number absent from the source is a DIFFERENT class: objectively
        # checkable, and wrong far more often than not. It is called out separately.
        if absent_tok or absent_num:
            hard.append(rec)
        (quoted if not absent_ng else composed).append(rec)

    print(f"{deck.name} — {len(lines)} shipped lines, traced against {len(read)} source file(s)")
    print(f"  quoted   (0 absent {a.n}-grams — near-verbatim from source) : {len(quoted)}")
    print(f"  composed (your own wording — where content defects live)    : {len(composed)}")

    if hard:
        print(f"\n  ⚠ {len(hard)} line(s) carry a NUMBER or IDENTIFIER absent from the source. "
              f"This class is exactly checkable and is usually a real defect:")
        for r in hard:
            bits = (r["absent_tokens"] + r["absent_numbers"])[:6]
            print(f"      s{r['slide']:>2}  {r['text'][:52]:<54} absent: {bits}")

    print(f"\n  COMPOSED lines — hand THESE to the content lens, not the whole deck:")
    for r in composed:
        print(f"      s{r['slide']:>2}  {r['text'][:66]}")

    print("\n  This is TRIAGE, not a verdict. A composed line is not wrong; it is simply not")
    print("  quoted, so nothing but a reader can confirm it. Quoted lines still need their")
    print("  numbers checked — see the ⚠ class above.")

    if a.json:
        Path(a.json).write_text(json.dumps(
            {"deck": str(deck), "sources": read, "n": a.n,
             "quoted": quoted, "composed": composed, "absent_hard": hard},
            ensure_ascii=False, indent=2))
        print(f"\n  [json] wrote {a.json}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
