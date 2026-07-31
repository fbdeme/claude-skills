#!/usr/bin/env python3
"""Codex-only visual-contract checks for fragile local slide relationships.

This is deliberately narrower than the shared layout lints. A contract is declared only
for high-risk manual relationships: a callout near a title, a component value beside a
neighbouring diagram, or an icon whose glyph carries meaning. The script validates
those relationships against the final PPTX and emits reviewer-ready crops from final
render PNGs. Claude Code's normal workflow does not invoke this script.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import sys
from pathlib import Path
from typing import Any


SCHEMA = "slide-maker-codex-visual-contract/v1"
RESULT_SCHEMA = "slide-maker-codex-visual-contract-result/v1"
EMU_PER_INCH = 914400.0


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return value


def bbox(shape: Any) -> tuple[float, float, float, float]:
    return (
        shape.left / EMU_PER_INCH,
        shape.top / EMU_PER_INCH,
        shape.width / EMU_PER_INCH,
        shape.height / EMU_PER_INCH,
    )


def rect_from(value: Any) -> tuple[float, float, float, float] | None:
    if not isinstance(value, list) or len(value) != 4 or not all(isinstance(v, (int, float)) for v in value):
        return None
    x, y, w, h = (float(v) for v in value)
    if w <= 0 or h <= 0:
        return None
    return x, y, w, h


def rect_gap(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> tuple[float, float, float]:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    dx = max(bx - (ax + aw), ax - (bx + bw), 0.0)
    dy = max(by - (ay + ah), ay - (by + bh), 0.0)
    return dx, dy, math.hypot(dx, dy)


def union(rects: list[tuple[float, float, float, float]]) -> tuple[float, float, float, float]:
    left = min(rect[0] for rect in rects)
    top = min(rect[1] for rect in rects)
    right = max(rect[0] + rect[2] for rect in rects)
    bottom = max(rect[1] + rect[3] for rect in rects)
    return left, top, right - left, bottom - top


def _leaf_shapes(shapes: Any) -> Any:
    """Every shape, descending through GROUPS.

    A top-level-only walk reported "text target not found" for any target inside a group, which on
    a designer-tool export or a redesign input is most of the deck — and the contract then failed
    for a reason that had nothing to do with the layout it was asserting. lint_deck's walker
    descends; this one has to agree with it or the two paths disagree about what is on the slide.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _leaf_shapes(shape.shapes)
        else:
            yield shape


def text_shapes(slide: Any, needle: str) -> list[tuple[float, float, float, float]]:
    return [
        bbox(shape)
        for shape in _leaf_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False) and (shape.text or "").strip() == needle
    ]


def target_rects(slide: Any, target: Any) -> tuple[list[tuple[float, float, float, float]], str | None]:
    if not isinstance(target, dict):
        return [], "target must be an object"
    if isinstance(target.get("text"), str) and target["text"].strip():
        matches = text_shapes(slide, target["text"].strip())
        occurrence = target.get("occurrence", 0)
        if not isinstance(occurrence, int) or occurrence < 0:
            return [], "target occurrence must be a non-negative integer"
        if occurrence >= len(matches):
            return [], f"text target not found: {target['text']!r} occurrence {occurrence}"
        return [matches[occurrence]], None
    rect = rect_from(target.get("rect"))
    if rect is None:
        return [], "target needs exact text or a positive [x, y, w, h] rect"
    return [rect], None


def crop_render(
    renders: Path | None,
    slide_number: int,
    rects: list[tuple[float, float, float, float]],
    output_dir: Path,
    zone_id: str,
    slide_width: float,
    slide_height: float,
) -> str | None:
    if renders is None:
        return None
    candidate = renders / f"slide{slide_number:02d}.png"
    if not candidate.is_file():
        return None
    try:
        from PIL import Image
        image = Image.open(candidate).convert("RGB")
    except Exception:
        return None
    x, y, w, h = union(rects)
    margin = 0.18
    x = max(0.0, x - margin)
    y = max(0.0, y - margin)
    w = min(slide_width - x, w + 2 * margin)
    h = min(slide_height - y, h + 2 * margin)
    px = int(round(x / slide_width * image.width))
    py = int(round(y / slide_height * image.height))
    pr = int(round((x + w) / slide_width * image.width))
    pb = int(round((y + h) / slide_height * image.height))
    output_dir.mkdir(parents=True, exist_ok=True)
    out = output_dir / f"{zone_id}.png"
    image.crop((px, py, pr, pb)).save(out)
    return str(out)


def evaluate(
    pptx: Path,
    manifest_path: Path,
    build_script: Path | None = None,
    renders: Path | None = None,
    crop_dir: Path | None = None,
) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for visual-contract checks") from exc

    manifest = load_json(manifest_path)
    if manifest.get("schema") != SCHEMA:
        raise ValueError(f"manifest.schema must be {SCHEMA!r}")
    presentation = Presentation(str(pptx))
    slide_width = presentation.slide_width / EMU_PER_INCH
    slide_height = presentation.slide_height / EMU_PER_INCH
    errors: list[str] = []
    zones_out: list[dict[str, Any]] = []
    seen_ids: set[str] = set()

    zones = manifest.get("zones", [])
    if not isinstance(zones, list):
        errors.append("manifest.zones must be a list")
        zones = []
    for index, zone in enumerate(zones, start=1):
        label = f"zones[{index}]"
        if not isinstance(zone, dict):
            errors.append(f"{label} must be an object")
            continue
        zone_id = zone.get("id")
        slide_number = zone.get("slide")
        minimum = zone.get("min_gap", 0.12)
        if not isinstance(zone_id, str) or not zone_id.strip() or zone_id in seen_ids:
            errors.append(f"{label}.id must be a unique non-empty string")
            continue
        seen_ids.add(zone_id)
        if not isinstance(slide_number, int) or not 1 <= slide_number <= len(presentation.slides):
            errors.append(f"{label}.slide is outside the final deck")
            continue
        if not isinstance(minimum, (int, float)) or minimum < 0:
            errors.append(f"{label}.min_gap must be a non-negative number")
            continue
        subject, problem = target_rects(presentation.slides[slide_number - 1], zone.get("subject"))
        if problem:
            errors.append(f"{zone_id}: {problem}")
            continue
        blockers = zone.get("clear_of")
        if not isinstance(blockers, list) or not blockers:
            errors.append(f"{zone_id}: clear_of must name at least one adjacent target")
            continue
        checks: list[dict[str, Any]] = []
        all_rects = list(subject)
        zone_pass = True
        for blocker in blockers:
            against, problem = target_rects(presentation.slides[slide_number - 1], blocker)
            if problem:
                zone_pass = False
                checks.append({"pass": False, "error": problem})
                continue
            all_rects.extend(against)
            for subject_rect in subject:
                for against_rect in against:
                    dx, dy, distance = rect_gap(subject_rect, against_rect)
                    passed = max(dx, dy) >= float(minimum)
                    if not passed:
                        zone_pass = False
                    checks.append({
                        "pass": passed,
                        "subject": [round(v, 3) for v in subject_rect],
                        "against": [round(v, 3) for v in against_rect],
                        "horizontal_gap": round(dx, 3),
                        "vertical_gap": round(dy, 3),
                        "distance": round(distance, 3),
                        "required_gap": float(minimum),
                    })
        crop = crop_render(renders, slide_number, all_rects, crop_dir or manifest_path.parent / "visual-crops", zone_id, slide_width, slide_height)
        zones_out.append({"id": zone_id, "slide": slide_number, "pass": zone_pass, "checks": checks, "crop": crop})
        if not zone_pass:
            errors.append(f"{zone_id}: a named local relationship is overlapping or closer than {minimum:.2f}in")

    build_text = build_script.read_text(encoding="utf-8") if build_script and build_script.is_file() else ""
    icons_out: list[dict[str, Any]] = []
    icons = manifest.get("icons", [])
    if not isinstance(icons, list):
        errors.append("manifest.icons must be a list")
        icons = []
    for index, icon in enumerate(icons, start=1):
        label = f"icons[{index}]"
        if not isinstance(icon, dict):
            errors.append(f"{label} must be an object")
            continue
        icon_id = icon.get("id")
        token = icon.get("build_token")
        meaning = icon.get("meaning")
        slide_number = icon.get("slide")
        if not isinstance(icon_id, str) or not icon_id.strip() or icon_id in seen_ids:
            errors.append(f"{label}.id must be a unique non-empty string")
            continue
        seen_ids.add(icon_id)
        valid = True
        reasons: list[str] = []
        if not isinstance(slide_number, int) or not 1 <= slide_number <= len(presentation.slides):
            valid = False; reasons.append("slide is outside the final deck")
        if not isinstance(meaning, str) or len(meaning.strip()) < 12:
            valid = False; reasons.append("meaning must explain the icon's semantic job")
        if not isinstance(token, str) or not token.strip():
            valid = False; reasons.append("build_token must identify the icon glyph")
        elif build_script is None:
            valid = False; reasons.append("build script is required to bind an icon token")
        elif token not in build_text:
            valid = False; reasons.append(f"build token {token!r} is absent from the build script")
        icons_out.append({"id": icon_id, "slide": slide_number, "pass": valid, "meaning": meaning, "build_token": token, "reasons": reasons})
        if not valid:
            errors.append(f"{icon_id}: " + "; ".join(reasons))

    return {
        "schema": RESULT_SCHEMA,
        "pptx": str(pptx),
        "pptx_sha256": sha256_file(pptx),
        "manifest": str(manifest_path),
        "manifest_sha256": sha256_file(manifest_path),
        "passed": not errors,
        "zones": zones_out,
        "icons": icons_out,
        "errors": errors,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Verify Codex-only visual-contract zones and icon semantics.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--build-script", type=Path)
    parser.add_argument("--renders", type=Path)
    parser.add_argument("--crops", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    args = parser.parse_args()
    try:
        result = evaluate(args.pptx, args.manifest, args.build_script, args.renders, args.crops)
    except Exception as exc:
        print(f"visual contract failed to run: {exc}", file=sys.stderr)
        return 2
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"visual contract: {'PASS' if result['passed'] else 'FAIL'} · {args.out}")
    for error in result["errors"]:
        print("  " + error)
    return 0 if result["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
